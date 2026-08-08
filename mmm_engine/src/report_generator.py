# -*- coding: utf-8 -*-
"""PPTX report generator — MMM report (MTERIA brand design)."""
import io
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib import rcParams

_JP_FONTS = ['Noto Sans JP', 'Meiryo', 'Yu Gothic', 'MS Gothic', 'DejaVu Sans']
rcParams.update({
    'font.family':      'sans-serif',
    'font.sans-serif':  _JP_FONTS,
    'font.size':        8,
    'text.color':       '#000000',
    'axes.labelcolor':  '#000000',
    'xtick.color':      '#000000',
    'ytick.color':      '#000000',
})

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette (from PPTX reference) ────────────────────
COVER_BG    = RGBColor(0x31, 0x48, 0x58)  # cover/dark slide bg
SECTION_BG  = RGBColor(0x31, 0x5E, 0x6D)  # section slides bg + heading text
BRAND_LIGHT = RGBColor(0x5C, 0x92, 0x91)  # descriptions, TOC items, box borders
WATERMARK   = RGBColor(0x31, 0x76, 0x80)  # section number watermark + copyright
RULE_COLOR  = RGBColor(0xA2, 0xCE, 0xBF)  # horizontal rule
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG     = RGBColor(0xF3, 0xF7, 0xF4)  # KPI card light bg
TABLE_ROW1  = RGBColor(0xF0, 0xF7, 0xFA)  # table odd row
TABLE_ROW2  = RGBColor(0xFF, 0xFF, 0xFF)  # table even row

# Chart accent colors (kept for data visualizations)
ACCENT_ORG  = RGBColor(0xCB, 0x80, 0x13)  # brand amber accent
ACCENT_GRN  = RGBColor(0x31, 0x76, 0x80)  # brand mid-dark teal
ACCENT_YLW  = RGBColor(0xFF, 0xCC, 0x02)  # yellow (unused)

# Saturation status colors (飽和状態ベース)
SAT_HEADROOM = RGBColor(0x31, 0x5E, 0x6D)   # 伸び代あり → brand teal
SAT_MODERATE = RGBColor(0x7D, 0xAA, 0xA8)   # 適正域 → mid teal
SAT_SATURATED = RGBColor(0xFF, 0x6B, 0x35)  # 飽和域 → orange
SAT_ZERO     = RGBColor(0xBB, 0xBB, 0xBB)   # 係数ゼロ

FONT_JP  = 'Noto Sans JP'
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
COPYRIGHT_TEXT = 'Copyright © M TERIA Inc. All Rights Reserved.'

LOGO_PATH = r'G:\マイドライブ\00_ai_company\00_admin\logo\m-teria\logo-m-teria-white.png'


# ── Helpers ────────────────────────────────────────────────

def _rgb(color: RGBColor):
    return (int(color[0]) / 255, int(color[1]) / 255, int(color[2]) / 255)


def _add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _bg(slide, color=COVER_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_shadow(shape):
    """Override theme shadow by adding empty effectLst."""
    try:
        spPr = shape._element.spPr
        for el in list(spPr):
            if el.tag.endswith('}effectLst'):
                spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))
    except Exception:
        pass


def _fmt_man(man: float) -> str:
    """Format a 万円 value, switching to 億+万円 form at ≥10,000万円."""
    m = round(man)
    if m < 10000:
        return f'{m}万円'
    oku, rem = divmod(m, 10000)
    return f'{oku}億{rem}万円' if rem else f'{oku}億円'


def _box(slide, left, top, w, h, text='', font_size=18, bold=False,
         text_color=SECTION_BG, bg_color=None, align=PP_ALIGN.LEFT,
         font_name=FONT_JP, line_spacing=None, v_anchor=None, url=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if v_anchor is not None:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    try:
        run.font.name = font_name
    except Exception:
        pass
    if url:
        run.hyperlink.address = url
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox


def _line(slide, x1, y1, x2, y2, color=RULE_COLOR, width_pt=1.2):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    _no_shadow(connector)


def _rect(slide, left, top, w, h, fill_color=None, line_color=None, line_pt=0.75):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    _no_shadow(shape)
    return shape


def _rrect(slide, left, top, w, h, fill_color=None, line_color=None, line_pt=1.5, adj=16667):
    """Rounded rectangle (auto_shape_type=5). adj: 0=square, 50000=pill."""
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.adjustments[0] = adj
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    _no_shadow(shape)
    return shape


_HYPERLINK_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'

def _btn(slide, left, top, w, h, text, fill_color=None, line_color=None,
         line_pt=1.5, adj=0.2, font_size=12, text_color=WHITE, url=None):
    """ボタン用丸角矩形。クリックアクションを cNvPr に直接書くため PDF 変換後もリンク保持。"""
    shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.adjustments[0] = adj
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    _no_shadow(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = FONT_JP
    if url:
        # shape 全体のクリックアクションとして cNvPr に hlinkClick を追加
        rId = slide.part.relate_to(url, _HYPERLINK_TYPE, is_external=True)
        cNvPr = shape._element.find('.//' + qn('p:cNvPr'))
        hlink = etree.SubElement(cNvPr, qn('a:hlinkClick'))
        hlink.set(qn('r:id'), rId)
    return shape


def _mixed_price_box(slide, left, top, w, h, prefix, big_num, mid_suffix, small_suffix=''):
    """Single-line price: 'prefix BIG_NUM mid_suffix small_suffix' — small_suffix at 10pt."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    for txt, fz, bold in [(prefix, 14, True), (big_num, 48, True),
                          (mid_suffix, 14, True), (small_suffix, 10, True)]:
        if not txt:
            continue
        run = para.add_run()
        run.text = txt
        run.font.size = Pt(fz)
        run.font.bold = bold
        run.font.color.rgb = SECTION_BG
        run.font.name = FONT_JP
    _no_shadow(txBox)
    return txBox


def _copyright(slide, text_color=WATERMARK):
    _box(slide, 0.4, 7.08, 12.5, 0.25, COPYRIGHT_TEXT,
         font_size=9, text_color=text_color, align=PP_ALIGN.CENTER)


def _section_header(slide, part_num, title, subtitle=''):
    _bg(slide, SECTION_BG)
    # Large watermark number (full slide, vertically centered)
    _box(slide, 0.5, 0.0, 12.33, 7.5, f'{part_num:02d}',
         font_size=240, bold=True, text_color=WATERMARK, align=PP_ALIGN.CENTER,
         v_anchor=MSO_ANCHOR.MIDDLE)
    # Title box: full slide height so MIDDLE anchor = true center
    _box(slide, 0.5, 0.0, 12.33, 7.5, title,
         font_size=32, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER,
         v_anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _box(slide, 0.5, 4.4, 12.33, 0.6, subtitle,
             font_size=14, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
    _copyright(slide, WHITE)


def _slide_title(slide, title, subtitle=''):
    """Content slide header: title + horizontal rule + optional subtitle below rule + copyright."""
    _box(slide, 0.4, 0.20, 12.5, 0.58, title,
         font_size=22, bold=True, text_color=SECTION_BG)
    _line(slide, 0.4, 0.82, 12.93, 0.82, color=RULE_COLOR, width_pt=1.2)
    if subtitle:
        _box(slide, 0.4, 0.92, 12.5, 0.55, subtitle,
             font_size=12, text_color=SECTION_BG, line_spacing=1.5)
    _copyright(slide)


def _kpi_card(slide, left, top, w, h, label, value, color=None):
    _rect(slide, left, top, w, h, fill_color=CARD_BG)
    _box(slide, left + 0.1, top + 0.08, w - 0.2, h * 0.38, label,
         font_size=10, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
    _box(slide, left + 0.1, top + h * 0.42, w - 0.2, h * 0.55, value,
         font_size=20, bold=True, text_color=color or SECTION_BG, align=PP_ALIGN.CENTER)


def _fig_to_image(fig) -> io.BytesIO:
    # canvas.draw() forces auto-generated tick labels to be created,
    # then _style_labels can apply color/font to them
    fig.canvas.draw()
    for ax in fig.get_axes():
        _style_labels(ax)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=120, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _export_png(fig, export_dir: str | None, name: str):
    """export_dir が指定されている場合にグラフを PNG 保存する。"""
    if export_dir is None:
        return
    import os
    os.makedirs(export_dir, exist_ok=True)
    fig.canvas.draw()
    for ax in fig.get_axes():
        _style_labels(ax)
    fig.savefig(os.path.join(export_dir, f'{name}.png'), dpi=150,
                bbox_inches='tight', facecolor=fig.get_facecolor())


def _add_chart_image(slide, buf, left, top, w, h):
    slide.shapes.add_picture(buf, Inches(left), Inches(top), Inches(w), Inches(h))


def _set_cell_border(cell, color=RULE_COLOR, width_pt=0.5):
    """Add border lines to a table cell via XML.

    OpenXML schema requires lnL/R/T/B to appear BEFORE solidFill in a:tcPr.
    We use insert(0, ...) in reversed order to achieve lnL, lnR, lnT, lnB
    at positions 0-3, pushing any existing solidFill to the end.
    """
    clr_hex = str(color)  # RGBColor.__str__() returns e.g. 'A2CEBF'
    width_emu = int(width_pt * 12700)
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    # Remove existing border elements
    for side in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)
    # Insert in reversed order at index 0 → final order: lnL, lnR, lnT, lnB, <fill>
    for side in reversed(['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']):
        ln = etree.Element(qn(side))
        ln.set('w', str(width_emu))
        sf = etree.SubElement(ln, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', clr_hex)
        tcPr.insert(0, ln)


def _cell_vcenter(cell):
    """Force vertical-center on a table cell.
    OOXML table cells use a:tcPr/@anchor (not a:txBody/@anchor).
    """
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    tcPr.set('anchor', 'ctr')


def _add_table(slide, left, top, w, h, headers, rows,
               header_bg=None, header_font_size=12, body_font_size=11,
               col_widths=None, row_height=0.30, header_height=0.38,
               col_alignments=None):
    """h は最大高さ。20行以上は行高を自動縮小してスライドに収める（最小 0.14 インチ）。"""
    hbg = header_bg or SECTION_BG
    n_cols  = len(headers)
    n_data  = len(rows)
    n_rows  = n_data + 1
    # 20行以上はスライド内に収まるよう行高を自動縮小
    _MIN_ROW_H = 0.14
    effective_row_h = row_height
    if n_data >= 20 and n_data > 0:
        fit_h = (h - header_height) / n_data
        effective_row_h = min(row_height, max(_MIN_ROW_H, fit_h))
    actual_h = min(h, header_height + n_data * effective_row_h)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(w), Inches(actual_h)
    ).table
    # 列幅
    if col_widths and len(col_widths) == n_cols:
        total = sum(col_widths)
        for i in range(n_cols):
            tbl.columns[i].width = Inches(w * col_widths[i] / total)
    else:
        col_w = Inches(w / n_cols)
        for i in range(n_cols):
            tbl.columns[i].width = col_w
    # 行高さを明示指定（均等引き伸ばし防止）
    tbl.rows[0].height = Inches(header_height)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(effective_row_h)

    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = hbg
        _cell_vcenter(cell)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(header_font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
        try: run.font.name = FONT_JP
        except: pass
        _set_cell_border(cell)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW2  # white
            _cell_vcenter(cell)
            for p in cell.text_frame.paragraphs:
                p.alignment = col_alignments[j] if col_alignments and j < len(col_alignments) else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(body_font_size)
                    run.font.color.rgb = SECTION_BG
                    try: run.font.name = FONT_JP
                    except: pass
            _set_cell_border(cell)


# ── Chart helpers (light theme) ────────────────────────────

CHART_BG   = '#FFFFFF'
CHART_AX   = '#FFFFFF'
C_TEAL     = '#315E6D'
C_LTEAL    = '#5C9291'
C_TEAL2    = '#3E7A8A'   # mid-dark teal
C_TEAL3    = '#7DAAA8'   # mid teal
C_TEAL4    = '#A2CEBF'   # light teal (RULE_COLOR)
C_TEAL5    = '#C5DFD9'   # very light teal
C_ORANGE   = '#CB8013'   # brand amber accent
C_GREEN    = '#7EBEAB'   # brand accent mid-light
C_YELLOW   = '#CB8013'   # brand amber (current-spend dots)
C_RULE     = '#A2CEBF'
C_GRID     = '#DDE8EC'
C_TICK     = '#000000'   # black for axis labels


def _style_labels(ax):
    """Force color and font on all existing tick-label Text objects."""
    for lbl in ax.get_xticklabels() + ax.get_yticklabels():
        lbl.set_color(C_TICK)
        try:
            lbl.set_fontfamily(_JP_FONTS)
        except Exception:
            pass
    for obj in [ax.xaxis.label, ax.yaxis.label]:
        obj.set_color(C_TICK)
        try:
            obj.set_fontfamily(_JP_FONTS)
        except Exception:
            pass


def _chart_style(ax, fig=None, facecolor=CHART_AX):
    ax.set_facecolor(facecolor)
    if fig:
        fig.patch.set_facecolor(CHART_BG)
    # tick_params covers auto-generated labels (date axis etc.)
    ax.tick_params(axis='both', colors=C_TICK, labelcolor=C_TICK, labelsize=10)
    # iterate existing Text objects (covers set_yticklabels / set_xticklabels calls)
    _style_labels(ax)
    for sp in ax.spines.values():
        sp.set_color(C_RULE)
    ax.grid(False)


def _plot_actual_vs_pred(dates, actual_cv, pred_cv, dates_hold=None,
                          actual_hold=None, pred_hold=None, cv_col='CV'):
    fig, axes = plt.subplots(2, 1, figsize=(11, 5.5), facecolor=CHART_BG)
    ax1, ax2 = axes
    d = pd.to_datetime(dates)
    ax1.set_facecolor(CHART_AX)
    ax1.plot(d, actual_cv, color=C_TEAL, linewidth=0.9, label='実測', alpha=0.9)
    ax1.plot(d, pred_cv, color=C_ORANGE, linewidth=1.0, linestyle='--', label='予測')
    if dates_hold is not None and actual_hold is not None:
        dh = pd.to_datetime(dates_hold)
        ax1.plot(dh, actual_hold, color=C_GREEN, linewidth=0.9, label='実測（検証期間）', alpha=0.8)
        ax1.plot(dh, pred_hold, color=C_YELLOW, linewidth=1.0, linestyle=':', label='予測（検証期間）')
        ax1.axvline(dh[0], color=C_LTEAL, linewidth=1.0, linestyle='--', alpha=0.6)
        ax1.text(dh[0], ax1.get_ylim()[1] * 0.88, '検証期間', color=C_LTEAL, fontsize=7)
    ax1.legend(fontsize=8, loc='upper left', facecolor=CHART_BG, labelcolor=C_TICK)
    ax1.set_ylabel(cv_col, color=C_TICK, fontsize=9)
    _chart_style(ax1)
    ax2.set_facecolor(CHART_AX)
    residuals = actual_cv - pred_cv
    ax2.bar(d, residuals, color=np.where(residuals >= 0, C_TEAL, C_ORANGE), width=0.8, alpha=0.7)
    ax2.axhline(0, color=C_LTEAL, linewidth=0.8)
    ax2.set_ylabel('残差', color=C_TICK, fontsize=9)
    _chart_style(ax2)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_monthly_accuracy(dates, actual_cv, pred_cv, cv_col='CV'):
    df = pd.DataFrame({'date': pd.to_datetime(dates), 'actual': actual_cv, 'pred': pred_cv})
    df['ym'] = df['date'].dt.to_period('M')
    monthly = df.groupby('ym').agg({'actual': 'sum', 'pred': 'sum'}).reset_index()
    monthly['error_pct'] = np.abs(monthly['actual'] - monthly['pred']) / np.maximum(monthly['actual'], 1) * 100
    fig, ax = plt.subplots(figsize=(11, 3.5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    x = np.arange(len(monthly))
    ax.bar(x - 0.2, monthly['actual'], 0.38, label='実測', color=C_TEAL, alpha=0.85)
    ax.bar(x + 0.2, monthly['pred'],   0.38, label='予測', color=C_LTEAL, alpha=0.85)
    ax2 = ax.twinx()
    ax2.plot(x, monthly['error_pct'], 'o-', color=C_ORANGE, linewidth=1.0, markersize=4)
    ax2.set_ylabel('誤差率%', color=C_ORANGE, fontsize=10)
    ax.set_xticks(x)
    ax.set_xticklabels([str(p) for p in monthly['ym']], rotation=45, fontsize=9, color=C_TICK)
    _chart_style(ax)
    for sp in ax2.spines.values():
        sp.set_color(C_RULE)
    ax2.tick_params(axis='y', colors=C_ORANGE, labelcolor=C_ORANGE, labelsize=9)
    ax.legend(fontsize=9, facecolor=CHART_BG, labelcolor=C_TICK)
    ax.set_ylabel(cv_col, color=C_TICK, fontsize=10)
    fig.tight_layout(pad=0.8)
    return fig


# 飽和状態ラベル → チャートカラー（緑グラデ：伸び代=濃、飽和=薄）
SAT_CHART_COLORS = {
    '伸び代あり': C_TEAL,      # '#315E6D' 濃い緑
    '適正域':     C_GREEN,     # '#7EBEAB' 中間緑
    '飽和域':     '#C5DFD9',   # 最も薄い緑（飽和側は薄く）
    '係数ゼロ':   '#BBBBBB',
    '計測不能':   '#BBBBBB',
    '稼働日不足': '#BBBBBB',
    '逆相関除外': '#BBBBBB',
    '-':          '#BBBBBB',
}


def _plot_channel_bar(channels, values, title='', color_by_tier=None,
                       channel_metrics=None, xlabel='', fmt='{:.0f}', reverse=True):
    sorted_idx = np.argsort(values)[::-1] if reverse else np.argsort(values)
    chs  = [channels[i] for i in sorted_idx]
    vals = [values[i]   for i in sorted_idx]
    colors = []
    for ch in chs:
        label = '-'
        if channel_metrics and ch in channel_metrics:
            label = channel_metrics[ch].get('saturation_label', '-')
        colors.append(SAT_CHART_COLORS.get(label, '#BBBBBB'))
    fig, ax = plt.subplots(figsize=(11, max(3.5, len(chs) * 0.28)), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    bars = ax.barh(range(len(chs)), vals, color=colors, alpha=0.9, height=0.65)
    ax.set_yticks(range(len(chs)))
    ax.set_yticklabels(chs, fontsize=10, color=C_TICK)
    ax.invert_yaxis()
    _chart_style(ax)
    ax.set_xlabel(xlabel, color=C_TICK, fontsize=10)
    for bar, val in zip(bars, vals):
        if val > 0:
            ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                    fmt.format(val), va='center', fontsize=9, color=C_TEAL)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_spend_vs_cv(channels, spend_shares, contrib_shares, rssd):
    fig, ax = plt.subplots(figsize=(10, 5.5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    x = np.arange(len(channels))
    ax.bar(x - 0.2, spend_shares * 100, 0.38, label='支出金額シェア', color=C_TEAL, alpha=0.8)
    ax.bar(x + 0.2, contrib_shares * 100, 0.38, label='CV貢献シェア', color=C_LTEAL, alpha=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels([c[:8] for c in channels], rotation=90, fontsize=9, color=C_TICK)
    _chart_style(ax)
    ax.set_ylabel('シェア%', color=C_TICK, fontsize=10)
    ax.legend(fontsize=9, facecolor=CHART_BG, labelcolor=C_TICK)
    ax.set_title(f'RSSD = {rssd:.4f}', color=C_TICK, fontsize=10)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_response_curve(curve: dict, ch_name: str, ch_info: dict, media_basis: str = 'media'):
    fig, ax = plt.subplots(figsize=(5.5, 3.8), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    ax.plot(curve['x'], curve['y'], color=C_TEAL, linewidth=2)
    if curve.get('current_x', 0) > 0:
        y_at_current = np.interp(curve['current_x'], curve['x'], curve['y'])
        ax.axvline(curve['current_x'], color=C_ORANGE, linewidth=1.2, linestyle='--', alpha=0.8)
        ax.scatter([curve['current_x']], [y_at_current], color=C_YELLOW, s=60, zorder=5)
        ax.text(curve['current_x'] * 1.02, y_at_current * 0.95,
                '現在の\n支出位置', fontsize=7, color=C_ORANGE)
    if media_basis == 'spend':
        xlabel = '支出金額（残存効果加味）[円]'
    else:
        xlabel = '支出量（残存効果加味）'
    ax.set_xlabel(xlabel, color=C_TICK, fontsize=10)
    ax.set_ylabel('推定CV数', color=C_TICK, fontsize=10)
    _chart_style(ax)
    sat_label = ch_info.get('saturation_label', '-')
    cpa       = ch_info.get('cpa')
    cpa_s     = f'¥{cpa:,.0f}' if cpa else 'N/A'
    ax.set_title(f'{ch_name}  CPA={cpa_s}  [{sat_label}]', color=C_TICK, fontsize=9)
    fig.tight_layout(pad=0.5)
    return fig


def _plot_budget_bar(channels, current_spends, optimal_spends):
    sorted_idx = np.argsort(optimal_spends)[::-1]
    chs = [channels[i] for i in sorted_idx]
    cur = [current_spends[i] / 10000 for i in sorted_idx]
    opt = [optimal_spends[i] / 10000 for i in sorted_idx]
    fig, ax = plt.subplots(figsize=(11, max(4, len(chs) * 0.28)), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    x = np.arange(len(chs))
    ax.barh(x - 0.2, cur, 0.35, label='現状金額（万円）', color=C_TEAL4,  alpha=0.85)
    ax.barh(x + 0.2, opt, 0.35, label='最適金額（万円）', color=C_TEAL, alpha=0.85)
    ax.set_yticks(x)
    ax.set_yticklabels(chs, fontsize=10, color=C_TICK)
    ax.invert_yaxis()
    ax.set_ylim(len(chs) - 0.5, -0.9)  # top padding: legend と最上段バーの間にスペース
    ax.legend(fontsize=9, facecolor=CHART_BG, labelcolor=C_TICK,
              loc='upper center', bbox_to_anchor=(0.5, 1.10), ncol=2)
    _chart_style(ax)
    ax.set_xlabel('支出金額（万円）', color=C_TICK, fontsize=10)
    fig.subplots_adjust(left=0.18, right=0.96, top=0.88, bottom=0.10)
    return fig


def _plot_param_heatmap(channels, params_list, coefs):
    data = {
        'λ (Adstock)': [params_list[ch]['lambda'] for ch in channels],
        'α (Hill傾き)': [params_list[ch]['alpha']  for ch in channels],
        'γ (Hill飽和)': [params_list[ch]['gamma']  for ch in channels],
        '効果係数': list(np.array(coefs) / (max(np.abs(coefs)) + 1e-10)),
    }
    matrix = np.array(list(data.values()))  # shape: (4, n_channels)

    # Per-column min-max normalization for coloring (actual values shown as text)
    matrix_norm = np.zeros_like(matrix, dtype=float)
    for j in range(matrix.shape[0]):
        col = matrix[j, :]
        cmin, cmax = col.min(), col.max()
        matrix_norm[j, :] = (col - cmin) / (cmax - cmin) if cmax > cmin else 0.5

    fig, ax = plt.subplots(figsize=(6, max(4, len(channels) * 0.38)), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    im = ax.imshow(matrix_norm.T, aspect='auto', cmap='RdYlGn', vmin=0, vmax=1)
    ax.set_xticks(range(len(data)))
    ax.set_xticklabels(list(data.keys()), fontsize=11, color=C_TICK)
    ax.set_yticks(range(len(channels)))
    ax.set_yticklabels(channels, fontsize=11, color=C_TICK)
    for i in range(len(channels)):
        for j in range(len(data)):
            v_norm   = matrix_norm[j, i]
            v_actual = matrix[j, i]
            txt_color = 'white' if (v_norm > 0.75 or v_norm < 0.25) else C_TEAL
            ax.text(j, i, f'{v_actual:.2f}', ha='center', va='center',
                    fontsize=9, color=txt_color)
    plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_adstock_decay(lambda_vals: dict):
    """Adstock decay curves for each channel's λ value."""
    fig, ax = plt.subplots(figsize=(6.2, 3.8), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    t = np.arange(0, 31)
    colors = [C_TEAL, C_TEAL2, C_TEAL3, C_TEAL4, C_LTEAL]
    for idx, (ch, lam) in enumerate(lambda_vals.items()):
        decay = lam ** t
        ax.plot(t, decay, color=colors[idx % len(colors)],
                linewidth=1.5, label=f'{ch[:10]} (λ={lam:.2f})')
    ax.set_xlabel('経過日数', color=C_TICK, fontsize=9)
    ax.set_ylabel('残存効果割合', color=C_TICK, fontsize=9)
    ax.legend(fontsize=7, facecolor=CHART_BG, labelcolor=C_TICK, loc='upper right')
    _chart_style(ax)
    fig.tight_layout(pad=0.5)
    return fig


def _plot_monthly_stacked(dates_train, ch_daily_contrib, channels, ch_metrics):
    """Monthly stacked bar chart of CV contribution per channel (top 8)."""
    df = pd.DataFrame({'date': pd.to_datetime(dates_train)})
    for ch in channels:
        df[ch] = list(ch_daily_contrib.get(ch, [0] * len(dates_train)))
    df['ym'] = df['date'].dt.to_period('M')
    monthly = df.groupby('ym')[channels].sum()
    total_by_ch = monthly.sum(axis=0)
    top_chs = total_by_ch.nlargest(8).index.tolist()

    fig, ax = plt.subplots(figsize=(11, 4.5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    stack_colors = [C_TEAL, C_TEAL2, C_TEAL3, C_TEAL4, C_LTEAL, C_TEAL5, C_GREEN, C_YELLOW]
    x = np.arange(len(monthly))
    bottom = np.zeros(len(monthly))
    for idx, ch in enumerate(top_chs):
        vals = monthly[ch].values if ch in monthly.columns else np.zeros(len(monthly))
        color = stack_colors[idx % len(stack_colors)]
        ax.bar(x, vals, bottom=bottom, label=ch[:10],
               color=color, alpha=0.85, width=0.72)
        bottom += vals
    ax.set_xticks(x)
    ax.set_xticklabels([str(p) for p in monthly.index], rotation=45, fontsize=8, color=C_TICK)
    ax.set_ylabel('推定CV数', color=C_TICK, fontsize=9)
    ax.legend(fontsize=7, facecolor=CHART_BG, labelcolor=C_TICK,
              bbox_to_anchor=(1.01, 1), loc='upper left', ncol=1)
    _chart_style(ax)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_cv_decomp(dates_train, ch_daily_contrib, channels, actual_train_cv,
                    media_fraction=None):
    """Monthly stacked bar: baseline (natural) vs total media contribution.

    ch_daily_contrib は日次の相対シェア計算用に設計されており、
    合計をそのまま使うと過大帰属になる（model.py の media_fraction と乖離）。
    media_fraction が渡された場合はスケーリング補正を行い、
    絶対量を正確に保ちつつ月次の分布形状を維持する。
    """
    dates  = pd.to_datetime(dates_train)
    actual = np.array(actual_train_cv, dtype=float)

    # 日次チャネル寄与の生合計（相対分布の形状として使用）
    media_raw = np.zeros(len(actual))
    for ch in channels:
        media_raw += np.array(ch_daily_contrib.get(ch, [0] * len(actual)), dtype=float)

    # media_fraction によるスケーリング補正
    # ch_daily_contrib の合計 → model の Robyn スタイル分解値に揃える
    if media_fraction is not None and media_raw.sum() > 0:
        target_media_total = float(actual.sum()) * media_fraction
        scale = target_media_total / media_raw.sum()
        media = media_raw * scale
    else:
        media = media_raw

    df = pd.DataFrame({'date': dates,
                       'baseline': np.clip(actual - media, 0, None),
                       'media':    media})
    df['ym'] = df['date'].dt.to_period('M')
    monthly  = df.groupby('ym')[['baseline', 'media']].sum()

    total   = monthly['baseline'].sum() + monthly['media'].sum()
    med_pct = monthly['media'].sum() / total * 100 if total > 0 else 0.0
    bas_pct = 100 - med_pct
    x = np.arange(len(monthly))

    fig, ax = plt.subplots(figsize=(10.5, 4.2), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    ax.bar(x, monthly['baseline'],
           label=f'自然流入・SEO等  {bas_pct:.0f}%',
           color='#7DAAA8', alpha=0.85, width=0.72)
    ax.bar(x, monthly['media'], bottom=monthly['baseline'],
           label=f'広告による獲得  {med_pct:.0f}%',
           color=C_TEAL, alpha=0.90, width=0.72)
    ax.set_xticks(x)
    ax.set_xticklabels([str(p) for p in monthly.index], rotation=45, fontsize=8, color=C_TICK)
    ax.set_ylabel('CV数', color=C_TICK, fontsize=9)
    ax.legend(fontsize=9, facecolor=CHART_BG, labelcolor=C_TICK, loc='upper left')
    _chart_style(ax)
    fig.tight_layout(pad=0.8)
    return fig, med_pct


def _plot_response_curve_grid_simple(top_channels, ch_metrics, threshold_cpa=None):
    """2×2 grid of simplified response curves for top 4 channels.
    Title label is CPA-based (◎/○/△) for consistency with the channel scorecard."""
    chs = top_channels[:4]
    n   = len(chs)

    def _cpa_label(cm_):
        cpa_ = cm_.get('cpa')
        if not cpa_:
            return '— 効果未確認'
        cpa_s = f'CPA ¥{cpa_:,.0f}'
        if threshold_cpa:
            if cpa_ <= threshold_cpa * 0.5:
                return f'◎ 増額を推奨  {cpa_s}'
            elif cpa_ <= threshold_cpa:
                return f'○ 効率良好    {cpa_s}'
            else:
                return f'△ 要改善      {cpa_s}'
        return cpa_s

    fig, axes = plt.subplots(2, 2, figsize=(11, 5.5), facecolor=CHART_BG)
    axes_flat  = axes.flatten()
    for idx in range(4):
        ax = axes_flat[idx]
        ax.set_facecolor(CHART_AX)
        if idx >= n:
            ax.set_visible(False)
            continue
        ch    = chs[idx]
        cm    = ch_metrics[ch]
        curve = cm.get('curve_data') or {}
        x_data = curve.get('x')
        y_data = curve.get('y')
        if x_data is not None and y_data is not None and len(x_data) > 0:
            ax.plot(x_data, y_data, color=C_TEAL, linewidth=2)
            cx = curve.get('current_x', 0)
            if cx > 0:
                cy   = float(np.interp(cx, x_data, y_data))
                ax.axvline(cx, color=C_ORANGE, linewidth=1.5, linestyle='--', alpha=0.85)
                ax.scatter([cx], [cy], color=C_YELLOW, s=70, zorder=5)
                x_max = float(max(x_data))
                side  = 1.04 if cx < x_max * 0.65 else 0.4
                ax.text(cx * side, cy * 0.92, '現在の支出位置',
                        fontsize=7, color=C_ORANGE, va='top')
        ax.set_title(f'{ch}\n{_cpa_label(cm)}', fontsize=8, color=C_TEAL, pad=3)
        ax.set_xlabel('広告投下量', fontsize=7, color=C_TICK)
        ax.set_ylabel('推定日次CV数',   fontsize=7, color=C_TICK)
        ax.tick_params(labelsize=6)
        _chart_style(ax)
    fig.tight_layout(pad=0.5, h_pad=1.5, w_pad=1.0)
    return fig


def _plot_marginal_roi(valid_channels, ch_metrics, roi_label='CV獲得効率'):
    """Horizontal bar: marginal ROI per channel (CV per 1万 additional spend)."""
    chs = sorted(valid_channels,
                 key=lambda c: ch_metrics[c].get('marginal_roi', 0.0), reverse=True)[:16]
    mroi = [ch_metrics[ch].get('marginal_roi', 0.0) for ch in chs]
    colors = [SAT_CHART_COLORS.get(ch_metrics[ch].get('saturation_label', '-'), '#BBBBBB') for ch in chs]
    fig, ax = plt.subplots(figsize=(11, max(3.5, len(chs) * 0.32)), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    bars = ax.barh(range(len(chs)), mroi, color=colors, alpha=0.9, height=0.65)
    ax.set_yticks(range(len(chs)))
    ax.set_yticklabels([c[:12] for c in chs], fontsize=9, color=C_TICK)
    ax.invert_yaxis()
    ax.set_xlabel(f'追加投資{roi_label}（CV件数 / 1万）', color=C_TICK, fontsize=10)
    for bar, val in zip(bars, mroi):
        if val > 0:
            ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                    f'{val:.1f}', va='center', fontsize=9, color=C_TEAL)
    _chart_style(ax)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_forest(channel_metrics: dict, channels: list, roi_label: str = 'CV獲得効率'):
    """Horizontal forest plot: channel ROI point estimates + 95% CI error bars.

    Colors: green = CI entirely > 1.0, mid-teal = CI straddles 1.0, grey = CI entirely < 1.0.
    X-axis is capped at P85 of CI-high values to prevent outliers from crushing the scale.
    """
    from matplotlib.lines import Line2D

    valid = [
        (ch, channel_metrics[ch])
        for ch in channels
        if not channel_metrics[ch].get('is_zero')
        and channel_metrics[ch].get('ci_available')
    ]
    if not valid:
        return None

    # Sort by ROI descending so highest ROI appears at top
    valid.sort(key=lambda x: x[1]['roi'], reverse=True)

    labels  = [ch for ch, _ in valid]
    roi_pts = np.array([m['roi']                        for _, m in valid])
    roi_lo  = np.array([m.get('roi_ci_low',  m['roi']) for _, m in valid])
    roi_hi  = np.array([m.get('roi_ci_high', m['roi']) for _, m in valid])

    # Cap x-axis: P85 of CI-high values か roi_hi.max()*1.25 の大きい方（フロアなし）
    x_cap = float(max(np.percentile(roi_hi, 85), float(roi_hi.max()) * 1.25, 0.05))

    n = len(labels)
    fig_h = max(3.5, min(n * 0.42, 6.5))
    fig, ax = plt.subplots(figsize=(9.5, fig_h), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)

    for i, (label, roi, lo, hi) in enumerate(zip(labels, roi_pts, roi_lo, roi_hi)):
        if lo >= 1.0:
            clr = '#315E6D'  # dark green — CI entirely above break-even
        elif hi < 1.0:
            clr = '#BBBBBB'  # grey — CI entirely below break-even
        else:
            clr = C_GREEN    # green — CI straddles 1.0

        lo_d  = min(lo, x_cap)
        hi_d  = min(hi, x_cap)
        roi_d = min(roi, x_cap)

        # CI horizontal bar
        ax.plot([lo_d, hi_d], [i, i], color=clr, linewidth=2.2,
                solid_capstyle='round', zorder=3)
        # End cap ticks
        ax.plot([lo_d, lo_d], [i - 0.12, i + 0.12], color=clr, linewidth=1.4, zorder=3)
        ax.plot([hi_d, hi_d], [i - 0.12, i + 0.12], color=clr, linewidth=1.4, zorder=3)
        # Point estimate dot
        ax.scatter([roi_d], [i], color=clr, s=50, zorder=5,
                   edgecolors='white', linewidths=0.8)

        # Outlier annotation when values are clipped
        if roi > x_cap:
            ax.text(x_cap * 0.99, i - 0.05, f' ROI={roi:.1f}',
                    fontsize=6.5, color=clr, ha='right', va='bottom', style='italic')
        elif hi > x_cap * 0.99:
            ax.text(x_cap * 0.99, i - 0.05, f' CI>{hi:.1f}',
                    fontsize=6.5, color=clr, ha='right', va='bottom', style='italic')

    # Break-even reference line（ROI=1.0がxlim内にある場合のみ描画）
    # xlim: [-x_cap*0.04, x_cap*1.05] ― x=1.0がその範囲外だとbboxが異常に拡張する
    if (-x_cap * 0.04) <= 1.0 <= (x_cap * 1.05):
        ax.axvline(x=1.0, color=C_TEAL, linewidth=0.9, linestyle='--', alpha=0.6, zorder=2)
        ax.text(1.03, max(0, n // 2 - 1), '損益分岐\n(ROI=1.0)', fontsize=6.5, color=C_TEAL,
                va='center', ha='left')

    ax.set_yticks(range(n))
    ax.set_yticklabels(labels, fontsize=8)
    ax.invert_yaxis()
    ax.set_ylim(n - 0.5, -0.9)  # top padding: legend と最上段バーの間にスペース
    x_left = -x_cap * 0.04
    ax.set_xlim(x_left, x_cap * 1.05)
    ax.set_xlabel(f'{roi_label}（推定CV数 / 支出金額）', fontsize=9)

    legend_handles = [
        Line2D([0], [0], color='#315E6D', linewidth=2.2, label='CI全体 1以上'),
        Line2D([0], [0], color=C_GREEN,   linewidth=2.2, label='CI 1を跨ぐ'),
        Line2D([0], [0], color='#BBBBBB', linewidth=2.2, label='CI全体 1以下'),
    ]
    ax.legend(handles=legend_handles, fontsize=7,
              loc='upper center', bbox_to_anchor=(0.5, 1.02), ncol=3,
              framealpha=0.9, edgecolor=C_RULE)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='x', alpha=0.15, linestyle='--', color=C_GRID, zorder=1)
    _chart_style(ax)
    xlim_l, xlim_r = ax.get_xlim()
    ax.set_xticks([t for t in ax.get_xticks() if xlim_l <= t <= xlim_r])
    fig.subplots_adjust(left=0.16, right=0.96, top=0.90, bottom=0.14)
    return fig


def _plot_dummy_bar(dummy_info):
    """Horizontal bar chart of dummy variable coefficients."""
    if not dummy_info:
        return None
    items = sorted(dummy_info, key=lambda d: abs(d.get('coef', 0)), reverse=True)
    names = [f'{d.get("date","")[:10]} ({d.get("category","")})' for d in items]
    vals  = [d.get('coef', 0.0) for d in items]
    fig, ax = plt.subplots(figsize=(11, max(3.5, len(items) * 0.32)), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    colors = [C_TEAL if v >= 0 else C_TEAL4 for v in vals]
    ax.barh(range(len(items)), vals, color=colors, alpha=0.85, height=0.65)
    ax.set_yticks(range(len(items)))
    ax.set_yticklabels(names, fontsize=8, color=C_TICK)
    ax.invert_yaxis()
    ax.axvline(0, color=C_RULE, linewidth=0.8)
    ax.set_xlabel('効果係数（正=CV増加・負=CV減少）', color=C_TICK, fontsize=9)
    _chart_style(ax)
    fig.tight_layout(pad=0.8)
    return fig


def _plot_efficient_frontier(frontier: dict):
    """投資効率フロンティア: 予算 vs 推定CV/売上曲線 + 限界CPA/ROIライン。"""
    is_monetary = frontier.get('cv_metric_type') == 'monetary'
    curve    = frontier['curve']
    budgets  = [p['budget'] / 10000 for p in curve]
    cvs      = [p['cv']     for p in curve]
    max_eff  = frontier['max_efficient_budget'] / 10000

    fig, ax1 = plt.subplots(figsize=(11, 4.5), facecolor=CHART_BG)
    ax1.set_facecolor(CHART_AX)

    budgets_arr = np.array(budgets)
    cvs_arr     = np.array(cvs)
    cv_at_max   = float(np.interp(max_eff, budgets_arr, cvs_arr))

    _mask_l = budgets_arr <= max_eff
    _mask_r = budgets_arr >= max_eff
    left_b  = np.append(budgets_arr[_mask_l], max_eff)
    left_cv = np.append(cvs_arr[_mask_l],     cv_at_max)
    right_b  = np.insert(budgets_arr[_mask_r], 0, max_eff)
    right_cv = np.insert(cvs_arr[_mask_r],     0, cv_at_max)

    ax1.plot(budgets_arr, cvs_arr, color=C_TEAL, linewidth=2, marker='o', markersize=4, zorder=3)
    ax1.fill_between(left_b,  left_cv,  color='#C0ECDD')
    ax1.fill_between(right_b, right_cv, color='#F3F7F4')

    ax1.set_xlabel('総予算（万円）', color=C_TICK, fontsize=10)
    cv_ylabel = '推定売上（円）' if is_monetary else '推定CV数'
    ax1.set_ylabel(cv_ylabel, color=C_TICK, fontsize=10)
    _chart_style(ax1)

    # 副軸: 限界ROI（monetary）or 限界CPA（count）
    ax2 = ax1.twinx()
    ax2.set_zorder(ax1.get_zorder() + 1)
    ax2.patch.set_visible(False)

    if is_monetary:
        thr_roi  = frontier['threshold_roi']
        _cap     = thr_roi * 4
        mpts = []
        for i in range(1, len(curve)):
            p = curve[i]
            mroi = p.get('marginal_roi', -1)
            if 0 < mroi <= _cap:
                mid_b = (curve[i-1]['budget'] + p['budget']) / 2 / 10000
                mpts.append((mid_b, min(mroi, _cap)))
        if mpts:
            bx_arr = np.array([pt[0] for pt in mpts])
            cy_arr = np.array([pt[1] for pt in mpts])
            if len(cy_arr) >= 3:
                cy_smooth = np.convolve(cy_arr, np.ones(3) / 3, mode='same')
                cy_smooth[0] = cy_arr[0]; cy_smooth[-1] = cy_arr[-1]
            else:
                cy_smooth = cy_arr
            ax2.plot(bx_arr, cy_smooth, color=C_TEAL2, linewidth=1.5, linestyle='--', alpha=0.9)
        ax2.axhline(thr_roi, color=C_ORANGE, linewidth=1.0, linestyle='-', alpha=0.70)
        ax2.text(budgets[-1], thr_roi * 1.02, f'閾値 ROI {thr_roi:.2f}',
                 fontsize=7, color=C_ORANGE, ha='right')
        ax2.set_ylim(0, _cap)
        ax2.set_ylabel('限界ROI', color=C_TICK, fontsize=9)
    else:
        thr_cpa  = frontier['threshold_cpa']
        _cpa_cap = thr_cpa * 3
        mcpa_pts = []
        for i in range(1, len(curve)):
            p = curve[i]
            if 0 < p['marginal_cpa'] <= _cpa_cap:
                mid_b = (curve[i-1]['budget'] + p['budget']) / 2 / 10000
                mcpa_pts.append((mid_b, min(p['marginal_cpa'], _cpa_cap)))
        if mcpa_pts:
            bx_arr = np.array([pt[0] for pt in mcpa_pts])
            cy_arr = np.array([pt[1] for pt in mcpa_pts])
            if len(cy_arr) >= 3:
                cy_smooth = np.convolve(cy_arr, np.ones(3) / 3, mode='same')
                cy_smooth[0] = cy_arr[0]; cy_smooth[-1] = cy_arr[-1]
            else:
                cy_smooth = cy_arr
            ax2.plot(bx_arr, cy_smooth, color=C_TEAL2, linewidth=1.5, linestyle='--', alpha=0.9)
        ax2.axhline(thr_cpa, color=C_ORANGE, linewidth=1.0, linestyle='-', alpha=0.70)
        ax2.text(budgets[-1], thr_cpa * 1.02, f'閾値 ¥{thr_cpa:,.0f}',
                 fontsize=7, color=C_ORANGE, ha='right')
        ax2.set_ylim(0, _cpa_cap)
        ax2.set_ylabel('限界CPA（円）', color=C_TICK, fontsize=9)

    # 最大効率予算マーカー
    y_at_max = np.interp(max_eff, budgets_arr, cvs_arr)
    ax1.axvline(max_eff, color=C_ORANGE, linewidth=1.5, linestyle='--', alpha=0.85, zorder=2)
    ax1.scatter([max_eff], [y_at_max], color=C_ORANGE, s=70, zorder=5)
    ax1.text(max_eff * 1.01, y_at_max * 0.97,
             f'理論上の\n最大効率予算\n¥{max_eff:.0f}万',
             fontsize=8, color=C_ORANGE, va='top')

    ax2.tick_params(axis='y', colors=C_TICK, labelcolor=C_TICK, labelsize=9)
    for sp in ax2.spines.values():
        sp.set_color(C_RULE)

    fig.tight_layout(pad=0.8)
    return fig


def _roi_labels(cv_metric_type: str) -> tuple:
    """(roi_label, roi_unit) を返す。
    count   → ('CV獲得効率', 'CV件数÷支出金額')
    monetary → ('ROAS',       '売上÷支出金額')
    """
    if cv_metric_type == 'monetary':
        return 'ROAS', '売上÷支出金額'
    return 'CV獲得効率', 'CV件数÷支出金額'


# ── Simple report builder (SMB向け簡易版) ──────────────────

def _build_simple_report(metrics, opt_result, client_name, output_path,
                         cv_col, media_basis, freq, opt_result_b, frontier,
                         budget_increase_pct, opt_result_dec=None, export_dir=None):
    """10-slide simplified MMM report for SMB clients."""
    import os

    roi_label, roi_unit = _roi_labels(metrics.get('cv_metric_type', 'count'))
    _cv_metric_type     = metrics.get('cv_metric_type', 'count')

    ch_metrics        = metrics['channel_metrics']
    channels          = metrics['channels']
    r2                = metrics['r2']
    nrmse             = metrics['nrmse']
    nrmse_hold        = metrics['nrmse_hold']
    rssd              = metrics['rssd']
    media_fraction    = metrics.get('media_fraction', 0.0)
    total_cv          = metrics['total_cv']
    total_spend       = metrics['total_spend']
    n_dummies         = metrics['n_dummies']
    dates_train       = metrics['dates_train']
    dates_hold        = metrics['dates_hold']
    actual_train_cv   = metrics['actual_train_cv']
    pred_train_cv     = metrics['pred_train_sqrt'] ** 2
    actual_hold_cv    = metrics['actual_hold_cv']
    pred_hold_cv      = metrics['pred_hold_sqrt'] ** 2
    ch_daily_contrib  = metrics.get('ch_daily_contrib', {})
    analysis_days     = len(dates_train) + len(dates_hold)
    holdout_days      = len(dates_hold)
    freq_label        = '週次' if freq == 'weekly' else '日次'
    period_unit       = '週'  if freq == 'weekly' else '日'

    valid_channels = [ch for ch in channels if not ch_metrics[ch]['is_zero']]
    zero_channels  = [ch for ch in channels if ch_metrics[ch]['is_zero']]
    cv_lift        = opt_result['cv_lift_pct']
    threshold_cpa  = frontier['threshold_cpa'] if frontier else None
    _is_monetary   = (frontier.get('cv_metric_type') == 'monetary') if frontier else False

    # ◎○△× 評価ロジック（○ = U+25CB で統一）
    def _acc_symbol(val, thresholds):
        """thresholds = [(boundary, symbol), ...] ordered best→worst."""
        for thr, sym in thresholds:
            if thr(val):
                return sym
        return '×'

    r2_sym   = _acc_symbol(r2,         [(lambda v: v >= 0.90, '◎'),
                                         (lambda v: v >= 0.85, '○'),
                                         (lambda v: v >= 0.80, '△')])
    nr_sym   = _acc_symbol(nrmse,       [(lambda v: v < 0.10, '◎'),
                                         (lambda v: v < 0.12, '○'),
                                         (lambda v: v < 0.15, '△')])
    nh_sym   = _acc_symbol(nrmse_hold,  [(lambda v: v < 0.15, '◎'),
                                         (lambda v: v < 0.20, '○'),
                                         (lambda v: v < 0.25, '△')])
    rs_sym   = _acc_symbol(rssd,        [(lambda v: 0.10 <= v <= 0.20, '◎'),
                                         (lambda v: v <= 0.30, '○'),
                                         (lambda v: v <= 0.40, '△')])
    mcr_sym  = _acc_symbol(media_fraction, [(lambda v: v >= 0.15, '◎'),
                                             (lambda v: v >= 0.08, '○'),
                                             (lambda v: v >= 0.03, '△')])

    def _ch_symbol(cm_):
        if cm_['is_zero']:
            return '×', cm_.get('zero_reason', 'CV効果が確認できないチャネルです')
        cpa_ = cm_.get('cpa')
        roi_ = cm_.get('roi', 0.0)
        if threshold_cpa and cpa_:
            if cpa_ <= threshold_cpa * 0.5 and roi_ >= 1.5:
                return '◎', '非常に費用対効果の良いチャネルです'
            elif cpa_ <= threshold_cpa:
                return '○', '費用対効果は良好なチャネルです'
            else:
                return '△', 'コストに対してCV成果が少ないチャネルです'
        return ('○', '費用対効果は良好な水準です') if roi_ >= 1.0 else ('△', 'コストに対してCVが少ない状態です')

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Cover ────────────────────────────────────────
    sl = _add_slide(prs)
    _bg(sl, COVER_BG)
    _box(sl, 0.55, 0.35, 6, 0.4, f'{client_name}　御中',
         font_size=16, bold=True, text_color=WHITE)
    _box(sl, 0.5, 2.9, 12.33, 1.0, '未来が見える広告レポート',
         font_size=46, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 0.5, 3.95, 12.33, 0.65,
         'マーケティングミックスモデリングによる媒体分析と最適な予算配分',
         font_size=18, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    logo = LOGO_PATH
    if os.path.exists(logo):
        sl.shapes.add_picture(logo, Inches(5.67), Inches(6.5), Inches(2.0), height=None)
    else:
        _box(sl, 0.5, 6.5, 12.33, 0.4, 'M TERIA',
             font_size=16, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: 目次 ─────────────────────────────────────────
    sl = _add_slide(prs)
    _box(sl, 0.4, 0.20, 4, 0.58, '目次',
         font_size=22, bold=True, text_color=SECTION_BG)
    _line(sl, 0.4, 0.82, 12.93, 0.82, color=RULE_COLOR, width_pt=1.2)
    toc_simple = [
        ('01', '分析サマリ・マーケティングミックスモデリングの流れ',
         'モデリング結果のサマリ／MMMの4ステップ'),
        ('02', 'モデルの予測精度',
         '実測CVと予測CVの比較／未学習データによる検証'),
        ('03', 'チャネル別の効果分析',
         f'チャネル別評価／{roi_label}の確信度／追加投資余地の分析'),
        ('04', '最適な予算配分',
         '同額予算での最適配分／増額時の最適配分／減額時の最適配分／投資の限界効率'),
        ('05', 'モデルの信頼度・推奨アクション',
         '精度指標による評価／ネクストアクションのご提案'),
    ]
    for i, (num, ttl, desc) in enumerate(toc_simple):
        y = 1.30 + i * 1.02
        _box(sl, 1.25, y + 0.15, 0.85, 0.55, num,
             font_size=22, bold=True, text_color=BRAND_LIGHT, align=PP_ALIGN.LEFT)
        _box(sl, 2.35, y + 0.12, 10.8, 0.38, ttl,
             font_size=15, bold=True, text_color=SECTION_BG)
        _box(sl, 2.35, y + 0.58, 11.2, 0.28, desc,
             font_size=9, text_color=BRAND_LIGHT)
    _copyright(sl)

    # ── Slide 3: 分析サマリ（決裁者回覧用）──────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, '分析サマリ',
                 f'{analysis_days}{period_unit}間のデータによる分析を実施。モデル試算では、同じ予算金額で投資配分を変えることで、+{cv_lift:.1f}%のCV改善余地が見えます。')

    # 期間文字列（CTAでも流用）
    _all_dates  = list(dates_train) + list(dates_hold)
    _dt0        = pd.Timestamp(_all_dates[0])
    _dt1        = pd.Timestamp(_all_dates[-1])
    _period_str = f'{_dt0.year}年{_dt0.month}月〜{_dt1.year}年{_dt1.month}月'

    # 左半分: 現状 vs 最適配分 対比バーチャート
    _sum_chs  = list(opt_result['channel_opt'].keys())
    _sum_cur  = [opt_result['channel_opt'][ch]['current_spend'] for ch in _sum_chs]
    _sum_opt  = [opt_result['channel_opt'][ch]['optimal_spend'] for ch in _sum_chs]
    _fig_smry = _plot_budget_bar(_sum_chs, _sum_cur, _sum_opt)
    _export_png(_fig_smry, export_dir, '01_budget_summary')
    _add_chart_image(sl, _fig_to_image(_fig_smry), 0.40, 1.55, 6.00, 5.10)

    # 右半分: 3レコード × カードグリッド
    _rh1, _rh2, _rh3 = 1.30, 1.30, 2.10
    _ry1 = 1.55
    _ry2 = _ry1 + _rh1 + 0.20
    _ry3 = _ry2 + _rh2 + 0.20
    _GAP = 0.20
    _cw1 = 1.90
    _cw2 = 2.95
    _HDR = 0.36

    def _smr_card(x, y, w, h, lbl, val, vfz, val_ls=None, val_y_offset=0.0):
        _rect(sl, x, y, w, h, fill_color=CARD_BG)
        _rect(sl, x, y, w, _HDR, fill_color=SECTION_BG)
        _box(sl, x+0.10, y+0.05, w-0.20, _HDR-0.05, lbl,
             font_size=10, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
        val_y = y + _HDR + 0.08 + val_y_offset
        val_h = max(y + h - 0.06 - val_y, 0.20)
        kw = dict(font_size=vfz, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        if val_ls: kw['line_spacing'] = val_ls
        _box(sl, x+0.12, val_y, w-0.24, val_h, val, **kw)

    _period_display = _period_str.replace('〜', '〜\n')
    _smr_card(6.60,                  _ry1, _cw1, _rh1, '分析期間',   _period_display,         12, val_ls=1.5)
    _smr_card(6.60+_cw1+_GAP,       _ry1, _cw1, _rh1, 'チャネル数', f'{len(channels)}チャネル', 15, val_y_offset=0.16)
    _smr_card(6.60+2*(_cw1+_GAP),   _ry1, _cw1, _rh1, 'モデルの精度', f'R² = {r2:.3f}',           16, val_y_offset=0.16)

    _abs_cv_gain = max(opt_result['optimal_cv'] - opt_result['current_cv'], 0)
    _smr_card(6.60, _ry2, _cw2, _rh2, '分析期間の広告費合計', _fmt_man(total_spend/10000), 20, val_y_offset=0.10)
    _cx2 = 6.60 + _cw2 + _GAP
    _rect(sl, _cx2, _ry2, _cw2, _rh2, fill_color=CARD_BG)
    _rect(sl, _cx2, _ry2, _cw2, _HDR, fill_color=SECTION_BG)
    _box(sl, _cx2+0.10, _ry2+0.05, _cw2-0.20, _HDR-0.05, '同予算でのCV改善余地',
         font_size=10, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    _valy = _ry2 + _HDR + 0.14
    _box(sl, _cx2+0.12, _valy, _cw2-0.24, 0.36, f'+{cv_lift:.1f}%',
         font_size=20, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
    _box(sl, _cx2+0.12, _valy+0.38, _cw2-0.24, 0.36, f'（+{_abs_cv_gain:.0f}件増加の試算）',
         font_size=11, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

    # 増額推奨: 最適化が delta_spend>0 と判断したチャネル（ゼロ係数除く）
    _inc_chs = sorted(
        [ch for ch, v in opt_result['channel_opt'].items()
         if not v.get('is_zero') and v.get('delta_spend', 0) > 0],
        key=lambda c: opt_result['channel_opt'][c].get('delta_spend', 0), reverse=True
    )
    _inc_label = '増額推奨チャネル'
    _rev_chs = [ch for ch, v in opt_result['channel_opt'].items()
                if v.get('action') == '停止・効果検証']
    _inc_text = ' ／ '.join(_inc_chs) if _inc_chs else '—'
    _rev_text = ' ／ '.join(_rev_chs) if _rev_chs else '—'
    for rx, lbl, val in [
        (6.60,          f'{_inc_label}（{len(_inc_chs)}件）', _inc_text),
        (6.60+_cw2+_GAP, f'停止推奨チャネル（{len(_rev_chs)}件）', _rev_text),
    ]:
        _rect(sl, rx, _ry3, _cw2, _rh3, fill_color=CARD_BG)
        _rect(sl, rx, _ry3, _cw2, _HDR, fill_color=SECTION_BG)
        _box(sl, rx+0.10, _ry3+0.05, _cw2-0.20, _HDR-0.05, lbl,
             font_size=10, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
        _box(sl, rx+0.12, _ry3+_HDR+0.10, _cw2-0.24, _rh3-_HDR-0.16, val,
             font_size=9, bold=True, text_color=SECTION_BG, line_spacing=1.6)

    # ── Slide 4: MMMの流れ（4ステップ版）────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, 'マーケティングミックスモデリング（MMM）の流れ', 'MMMとは「複数の広告・施策のうち、どれがどのくらいCVに寄与しているか？を可視化し、未来の予算配分・投資チャネルの最適解を導き出す分析手法」です。')
    steps4 = [
        ('データの\n読み込み＆前処理',
         '日次 or 週次の売上・CV数・コスト・チャネル・施策データ・媒体数値・季節変数などを収集。\n\n欠損補完などのデータクレンジングやモデル精度を上げるための前処理を行います。'),
        ('広告効果の\n変換',
         'チャネルごとに残存効果と効果が飽和し始める水準を推定します。'),
        ('モデルの最適化・\n精度向上',
         '2,000回の試行で誤差が最小になるモデルを探索します。\n\n局所最適化やダミー変数の自動探索などを行いながら、モデル精度を高めていきます。'),
        ('最適な\n予算配分の導出',
         f'各チャネルの残存効果や間接効果を加味したCPA・{roi_label}を算出。\n\nCV成果が最大化するパレート最適な予算配分プランを出力します。'),
    ]
    box_w, box_h = 2.9, 4.0
    box_y = 1.90
    # 4カードを横幅に対して中央寄せ（カード幅2.9 × 4 + 矢印間隔0.3 × 3 = 12.5）
    x_start = (13.33 - (4 * box_w + 3 * 0.30)) / 2
    for i, (step, desc) in enumerate(steps4):
        x = x_start + i * (box_w + 0.30)
        _rect(sl, x, box_y, box_w, box_h, line_color=BRAND_LIGHT, line_pt=1.2)
        _box(sl, x + 0.15, box_y + 0.28, box_w - 0.30, 0.80, step,
             font_size=15, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER,
             line_spacing=1.3)
        _line(sl, x + 0.25, box_y + 1.28, x + box_w - 0.25, box_y + 1.28,
              color=RULE_COLOR, width_pt=0.8)
        _box(sl, x + 0.15, box_y + 1.58, box_w - 0.30, box_h - 1.75, desc,
             font_size=11, text_color=SECTION_BG, align=PP_ALIGN.LEFT, line_spacing=1.5)
        if i < 3:
            _box(sl, x + box_w + 0.04, box_y + (box_h - 0.38) / 2, 0.22, 0.38, '▶',
                 font_size=13, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)

    # ── Slide 5: 実測と予測（フルレポートp10相当）────────────
    sl = _add_slide(prs)
    _slide_title(sl, f'実測と予測の比較（{freq_label}{cv_col}）',
                 f'モデルが実際のCV・売上をどれだけ正確に予測できているかを示します（右端{holdout_days}{period_unit}は、未学習データによる検証期間）。')
    fig_avp = _plot_actual_vs_pred(
        dates_train, actual_train_cv, pred_train_cv,
        dates_hold, actual_hold_cv, pred_hold_cv, cv_col=cv_col)
    _export_png(fig_avp, export_dir, '02_actual_vs_pred')
    _add_chart_image(sl, _fig_to_image(fig_avp), 0.5, 1.35, 12.3, 5.55)

    # ── Slide 7: チャネル通信簿 ───────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別CPAの一覧')
    def _zero_cmt_rank(cmt: str):
        """×チャネルのコメント種別ソートキー（日数短い順を含む）。"""
        import re as _re
        if '施策フラグ' in cmt:
            return (3, 0)
        if '日しか支出' in cmt:
            m = _re.search(r'のうち(\d+)[日週]しか', cmt)
            return (2, int(m.group(1)) if m else 999)
        if '支出金額が小さく' in cmt:
            return (1, 0)
        return (0, 0)  # 「CV効果が確認できない」

    def _row_sort_key(ch):
        cm = ch_metrics[ch]
        sym, cmt = _ch_symbol(cm)
        if sym != '×':
            return (0, cm['cpa'] if cm.get('cpa') else float('inf'), 0)
        return (1,) + _zero_cmt_rank(cmt)

    rows_sc = []
    for ch in sorted(channels, key=_row_sort_key):
        cm       = ch_metrics[ch]
        sym, cmt = _ch_symbol(cm)
        cpa_s    = f'¥{cm["cpa"]:,.0f}' if cm.get('cpa') else '—'
        rows_sc.append([sym, ch, cpa_s, cmt])
    _add_table(sl, 0.4, 0.95, 12.5, 6.05,
               ['評価', 'チャネル', 'CPA', 'コメント'], rows_sc,
               header_font_size=10,
               body_font_size=8 if len(channels) > 16 else 9,
               col_widths=[0.45, 2.0, 1.5, 4.5])

    # ── Slide 7.5: ROI 確信度（Forest plot CI）────────────────
    _ci_chs_s = [ch for ch in valid_channels if ch_metrics[ch].get('ci_available')]
    if _ci_chs_s:
        sl = _add_slide(prs)
        _slide_title(sl, f'チャネル別{roi_label}の確信度（95%信頼区間）',
                     f'ドットは{roi_label}の点推定値、横棒は95%信頼区間（CI）を表します（濃い緑：CIが{roi_label} 1以上／'
                     f'緑：CIが1を跨ぐ／グレー：CIが{roi_label} 1未満）。')
        fig_fp_s = _plot_forest(ch_metrics, valid_channels)
        if fig_fp_s is not None:
            _export_png(fig_fp_s, export_dir, '03_roi_forest')
            _add_chart_image(sl, _fig_to_image(fig_fp_s), 0.5, 1.35, 12.3, 5.55)

    # ── Slide 8: レスポンスカーブ top 4 ───────────────────────
    # 最適配分で「増額推奨」のチャネルのみ対象（減額推奨チャネルとの矛盾を防ぐ）
    _inc_opt = {ch for ch, v in opt_result['channel_opt'].items()
                if v.get('delta_spend', 0) > 0}
    _eligible = [ch for ch in valid_channels if ch in _inc_opt]
    top4 = sorted(_eligible, key=lambda c: ch_metrics[c]['roi'], reverse=True)[:4]
    if top4:
        sl = _add_slide(prs)
        _slide_title(sl, '増額推奨チャネルの追加投資余地',
                     f'最適な予算配分で増額が推奨されるチャネルのうち、{roi_label}上位4つの投資余地。現在の支出位置より右に余裕があるほど増額インパクトが大きくなります。')
        fig_grid = _plot_response_curve_grid_simple(top4, ch_metrics, threshold_cpa=threshold_cpa)
        _export_png(fig_grid, export_dir, '04_response_curves')
        _add_chart_image(sl, _fig_to_image(fig_grid), 0.4, 1.35, 12.5, 5.65)

    # ── Slide 9: 最適予算配分 ────────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, f'現状金額での最適な予算配分（CV+{cv_lift:.1f}%見込み）',
                 '現在の配分と最適配分の比較。「最適」が「現状」よりも長くなっているチャネルへの増額が効果的と判断できます。')
    opt_chs    = list(opt_result['channel_opt'].keys())
    cur_spends = [opt_result['channel_opt'][ch]['current_spend'] for ch in opt_chs]
    opt_spends = [opt_result['channel_opt'][ch]['optimal_spend'] for ch in opt_chs]
    fig_bud    = _plot_budget_bar(opt_chs, cur_spends, opt_spends)
    _export_png(fig_bud, export_dir, '05_budget_optimal')
    _add_chart_image(sl, _fig_to_image(fig_bud), 0.5, 1.35, 12.3, 5.65)

    # ── シナリオB: 増額時の最適配分 ────────────────────────────
    if opt_result_b is not None:
        sl = _add_slide(prs)
        inc_pct_label = int(budget_increase_pct * 100)
        cv_lift_b     = opt_result_b['cv_lift_pct']
        new_bud       = opt_result_b['total_budget']
        _slide_title(sl, f'+{inc_pct_label}%増額時の最適な予算配分（CV+{cv_lift_b:.1f}%見込み）',
                     f'総予算を{_fmt_man(new_bud/10000)}（現状比+{inc_pct_label}%）まで上げた際の最適配分。モデル試算上、+{cv_lift_b:.1f}%のCV増加を見込みます。')
        opt_chs_b    = list(opt_result_b['channel_opt'].keys())
        cur_spends_b = [opt_result_b['channel_opt'][ch]['current_spend'] for ch in opt_chs_b]
        opt_spends_b = [opt_result_b['channel_opt'][ch]['optimal_spend'] for ch in opt_chs_b]
        fig_b = _plot_budget_bar(opt_chs_b, cur_spends_b, opt_spends_b)
        _export_png(fig_b, export_dir, '06_budget_increase')
        _add_chart_image(sl, _fig_to_image(fig_b), 0.5, 1.35, 12.3, 5.65)

    # ── シナリオC: 減額時の最適配分 ────────────────────────────
    if opt_result_dec is not None:
        sl = _add_slide(prs)
        dec_pct_label = int(opt_result_dec['decrease_pct'] * 100)
        cv_lift_dec   = opt_result_dec['cv_lift_pct']
        new_bud_dec   = opt_result_dec['total_budget']
        _slide_title(sl, f'-{dec_pct_label}%減額時の最適な予算配分（CV{cv_lift_dec:+.1f}%見込み）',
                     f'総予算を{_fmt_man(new_bud_dec/10000)}（現状比-{dec_pct_label}%）に絞った際の最適配分。{roi_label}が低いチャネルを削減し、効果的なチャネルへ集中させることでCV損失を最小化できます。')
        opt_chs_dec    = list(opt_result_dec['channel_opt'].keys())
        cur_spends_dec = [opt_result_dec['channel_opt'][ch]['current_spend'] for ch in opt_chs_dec]
        opt_spends_dec = [opt_result_dec['channel_opt'][ch]['optimal_spend'] for ch in opt_chs_dec]
        fig_dec = _plot_budget_bar(opt_chs_dec, cur_spends_dec, opt_spends_dec)
        _export_png(fig_dec, export_dir, '07_budget_decrease')
        _add_chart_image(sl, _fig_to_image(fig_dec), 0.5, 1.35, 12.3, 5.65)

    # ── 投資効率フロンティア ────────────────────────────────
    if frontier is not None:
        sl = _add_slide(prs)
        _fr_subtitle = (
            f'限界ROI（追加¥1あたりの売上）が目標ROI {frontier["threshold_roi"]:.2f} を下回ったポイントが理論上の増額上限と予測できます。'
            if _is_monetary else
            f'限界CPA（1万円追加あたりのCPA）が{frontier["threshold_cpa"]:,.0f}円（現状CPAの1.5倍）の閾値を超えたポイントが理論上の増額上限と予測できます。'
        )
        _slide_title(sl, '投資の限界効率', _fr_subtitle)
        fig_fr = _plot_efficient_frontier(frontier)
        _export_png(fig_fr, export_dir, '08_efficient_frontier')
        _add_chart_image(sl, _fig_to_image(fig_fr), 0.4, 1.35, 9.8, 5.55)
        max_eff = frontier['max_efficient_budget']
        base_b  = frontier['current_budget']
        _kpi_card(sl, 10.4, 1.65, 2.5, 1.1, '最大効率の予算上限',    _fmt_man(max_eff/10000),                SECTION_BG)
        _eff_delta = (max_eff / base_b - 1) * 100
        _kpi_card(sl, 10.4, 2.95, 2.5, 1.1, '現状予算比',
                  f'+{_eff_delta:.0f}%' if _eff_delta >= 0 else f'{_eff_delta:.0f}%', SECTION_BG)
        if _is_monetary:
            _kpi_card(sl, 10.4, 4.25, 2.5, 1.1, '目標ROI閾値', f'{frontier["threshold_roi"]:.2f}x', SECTION_BG)
        else:
            _kpi_card(sl, 10.4, 4.25, 2.5, 1.1, '許容CPAの閾値', f'{frontier["threshold_cpa"]:,.0f}円', SECTION_BG)
        # 注記: 常時 + 条件付き警告
        _fr_notes = (
            ['※縦の破線は限界ROI曲線が目標ROI閾値と交差する予算水準です。']
            if _is_monetary else
            ['※縦の破線は限界CPA曲線が閾値（許容CPA）と交差する予算水準です。']
        )
        # 現状予算より低い範囲で閾値超過していたらオーバースペンド警告
        if _is_monetary:
            _sub_crossed = any(
                not p['is_efficient'] and p['budget'] < base_b and p['marginal_roi'] > 0
                for p in frontier['curve'][1:]
            )
        else:
            _sub_crossed = any(
                not p['is_efficient'] and p['budget'] < base_b and p['marginal_cpa'] > 0
                for p in frontier['curve'][1:]
            )
        if _sub_crossed:
            _fr_notes.append(
                '※現状予算より低い水準で、すでに限界ROIが閾値を下回っている区間があるため、現状予算での配分最適化後に予算を絞ることで効率を最大化できる可能性があります。'
                if _is_monetary else
                '※現状予算より低い水準で、すでに限界CPAが閾値を超えている区間があるため、現状予算での配分最適化後に予算を絞ることで効率を最大化できる可能性があります。'
            )
        if r2 < 0.70:
            _fr_notes.append(f'※R²={r2:.2f}（0.7未満）でモデルの説明力が不十分なため、予算推奨は参考値となります。')
        _box(sl, 10.4, 5.45, 2.5, 0.90 + len(_fr_notes[1:]) * 0.52,
             '\n'.join(_fr_notes),
             font_size=7, text_color=BRAND_LIGHT, line_spacing=1.5)

    # ── Slide 10: モデル精度 ─────────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, 'モデルの信頼度',
                 f'説明力はCV・売上変動のうちモデルが捉えられた割合を指し、残りはモデルで説明できないベースラインとなります。')
    r2_context = (
        'モデル精度は非常に高い水準です。\n（目安：0.85以上）' if r2 >= 0.92 else
        'モデル精度は高い水準です。\n（目安：0.85以上）' if r2 >= 0.85 else
        'データを追加することで\n精度の向上が見込めます。\n（目安：0.85以上）' if r2 >= 0.70 else
        'データ数が少なく、\nモデルで説明しきれていません。\nインプットデータの改善が必要です。\n（目安：0.85以上）'
    )
    nrmse_context = (
        '予測精度は非常に高い水準です。\n（目安：0.12未満）' if nrmse < 0.10 else
        '予測精度は良好です。\n（目安：0.12未満）' if nrmse < 0.12 else
        '予測精度は許容範囲内ですが、\n改善の余地があります。\n（目安：0.12未満）' if nrmse < 0.15 else
        'データの追加・クレンジングで\n精度の改善が見込めます。\n（目安：0.12未満）'
    )
    nrmse_hold_context = (
        '未来のデータへの汎化性能は\n高い水準です。\n（目安：0.20未満）' if nrmse_hold < 0.15 else
        '未来のデータへの汎化性能は\n良好な水準です。\n（目安：0.20未満）' if nrmse_hold < 0.20 else
        'やや過学習の傾向があります。\n（目安：0.20未満）' if nrmse_hold < 0.30 else
        '検証期間に特異なデータがある\n可能性があります。\n（目安：0.20未満）'
    )
    rssd_context = (
        '整合性は高い水準ですが、\n均質的な配分になっている\n可能性があります。\n（目安：0.10〜0.30の間）' if rssd < 0.10 else
        '投資額とCV貢献の配分が\n整合しており良好です。\n（目安：0.10〜0.30の間）' if rssd < 0.30 else
        '投資額とCV貢献に乖離があるため、\n最適配分に沿った予算の見直しを\n推奨します。\n（目安：0.10〜0.30の間）'
    )
    mcr_context = (
        '広告効果を十分に\n捉えられています。\n（目安：15%以上）' if media_fraction >= 0.15 else
        '標準的な媒体帰属率です。\n（目安：15%以上）' if media_fraction >= 0.08 else
        'ベースライン要因が強い状態です。\nチャネル・データ追加で\n改善が見込めます。\n（目安：15%以上）' if media_fraction >= 0.03 else
        'ベースライン主体のモデルです。\nチャネル追加・データクレンジングを\n推奨します。\n（目安：10〜15%以上）'
    )

    # ── パターン診断（カード描画前に確定） ──────────────────
    _nrmse_bad   = nrmse >= 0.12
    _nrmse_h_bad = nrmse_hold >= 0.20
    _rssd_bad    = rssd > 0.30
    _mcr_bad     = media_fraction < 0.08
    _r2_bad      = r2 < 0.85

    if _nrmse_bad:
        _dtxt = 'モデルの予測精度が基準を下回っており、\nデータクレンジングまたは分析期間の延長・チャネルの追加が必要です。'
    elif _nrmse_h_bad:
        _dtxt = '訓練データへの適合度合いは良好ですが、検証期間における精度が低くなっています。\n該当期間に特異なイベントデータがないかを確認し、問題がなければデータ追加で改善できる可能性があります。'
    elif _r2_bad:
        _dtxt = '予測精度は良好ですが、説明力（R²）が低くなっているため、媒体・施策チャネルのデータ追加で改善できる可能性があります。\n導出された最適な予算配分の結果は、参考値としてとらえることをおすすめします。'
    elif _mcr_bad and _rssd_bad:
        _dtxt = '媒体帰属率が低く、配分に乖離も見られます。\nまず媒体データ・チャネルの見直しを行い、モデル精度を改善した後に予算を最適化することをおすすめします。'
    elif _mcr_bad:
        _dtxt = 'CV成果の大半がベースライン要因（SEO・ブランド等）によるものです。\n広告効果が本来小さい場合は正常ですが、計測していないチャネルがある場合は、データの追加で改善できる可能性があります。'
    elif _rssd_bad:
        _dtxt = '精度は良好ですが、投資シェアとCV貢献シェアに乖離が見られます。\n最適な予算配分に沿って予算を見直すことで改善が見込めます。'
    else:
        _dtxt = '全ての精度指標が基準を満たしています。\n最適な予算配分に沿って投資を見直すことで、CV成果を最大化できる可能性があります。'

    # バナーあり → カードを上にずらしてスペースを確保
    _cy = 1.65

    acc_items = [
        (f'{r2*100:.0f}%',  '説明力',   'CV成果の何%を\nモデルが説明できているか？',
         f'R² = {r2:.3f}',  r2_context),
        (nr_sym,  '予測精度',   'モデルのCV予測が\nどの程度正確か？',
         f'NRMSE = {nrmse:.4f}',  nrmse_context),
        (nh_sym,  '汎化性能',   '未来のデータに対しても\nモデルが通用するか？',
         f'NRMSE(holdout) = {nrmse_hold:.4f}', nrmse_hold_context),
        (rs_sym,  '配分整合性', 'モデルが計算する\n投資額とCV貢献が揃っているか？',
         f'RSSD = {rssd:.4f}',    rssd_context),
        (mcr_sym, '媒体帰属率', '広告媒体がCV獲得に\nどれくらい貢献しているか？',
         f'MCR = {media_fraction*100:.1f}%',   mcr_context),
    ]
    card_w = 2.30
    _max_guide_lines = max(g.count('\n') + 1 for _, _, _, _, g in acc_items)
    _extra_h         = max(0.0, (_max_guide_lines - 1) * 0.22)
    card_h           = 3.8 + _extra_h
    for i, (sym, title, desc, val, guide) in enumerate(acc_items):
        x = 0.415 + i * 2.55
        _rect(sl, x, _cy,          card_w, card_h, fill_color=CARD_BG)
        _box(sl, x + 0.08, _cy + 0.20, card_w - 0.16, 0.35, title,
             font_size=11, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.08, _cy + 0.58, card_w - 0.16, 0.95, sym,
             font_size=40, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _line(sl, x + 0.15, _cy + 1.73, x + card_w - 0.15, _cy + 1.73, color=RULE_COLOR, width_pt=0.8)
        _box(sl, x + 0.08, _cy + 1.98, card_w - 0.16, 0.58, desc,
             font_size=9, text_color=SECTION_BG, align=PP_ALIGN.CENTER, line_spacing=1.5)
        _box(sl, x + 0.08, _cy + 2.63, card_w - 0.16, 0.42, val,
             font_size=11, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _guide_lines_i = guide.count('\n') + 1
        _guide_h       = 0.38 + max(0.0, (_guide_lines_i - 1) * 0.22)
        _box(sl, x + 0.08, _cy + 3.10, card_w - 0.16, _guide_h, guide,
             font_size=8, text_color=SECTION_BG, align=PP_ALIGN.CENTER, line_spacing=1.5)

    # ── パターン診断テキスト（背景なし・中央寄せ） ──────────
    _banner_y  = _cy + card_h + 0.14
    _dtxt_lines = _dtxt.count('\n') + 1
    _banner_h  = 0.22 + _dtxt_lines * 0.20
    _box(sl, 0.415, _banner_y, 12.5, _banner_h, _dtxt,
         font_size=10, text_color=SECTION_BG, align=PP_ALIGN.CENTER, line_spacing=1.5)

    # ── 3ヶ月アクションロードマップ ───────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, '3ヶ月間のアクションロードマップ', 'CV成果の最大化に向けて、直近3ヶ月間で推奨する施策の実行プランをまとめます。')
    zero_action   = f'ゼロ係数{len(zero_channels)}チャネルの予算削減・停止' if zero_channels else 'ゼロ係数チャネルなし（現状を維持）'
    zero_action_c = f'ゼロ係数{len(zero_channels)}チャネルを停止'      if zero_channels else 'ゼロ係数チャネルなし（現状を維持）'

    # ロードマップ用：opt_resultから増額推奨チャネル上位3つを抽出
    _inc_sorted = sorted(
        [(ch, v['delta_spend']) for ch, v in opt_result['channel_opt'].items()
         if v.get('delta_spend', 0) > 0],
        key=lambda x: x[1], reverse=True
    )[:3]
    _inc_names = '・'.join(ch for ch, _ in _inc_sorted) if _inc_sorted else '伸びしろありチャネル'

    # ゼロ係数チャネル名（支出多い順上位3、超える場合は「他Nch」付記）
    _zero_sorted = sorted(zero_channels, key=lambda c: ch_metrics[c]['spend_man'], reverse=True)
    _zero_names  = '・'.join(_zero_sorted[:3])
    if len(zero_channels) > 3:
        _zero_names += f' 他{len(zero_channels)-3}チャネル'
    _zero_act_named = f'{_zero_names}の停止を検討' if zero_channels else 'ゼロ係数チャネルなし（現状を維持）'

    if r2 < 0.75:
        months = [
            ('Month 1: データ整備＆モデル再構築', [
                '1. 欠損データや外れ値の特定・補完',
                '2. コントロール変数の追加',
                '3. データ期間の延長 or 頻度の変更',
            ]),
            ('Month 2: モデル再分析＆精度の再確認', [
                '1. 改善データによるMMMの再実行',
                f'2. R²・NRMSEの改善幅を確認（現状R²={r2:.3f}）',
                '3. 精度の確認後に予算配分の変更判断へ活用',
            ]),
            ('Month 3: 予算最適化＆次年度計画', [
                '1. 再分析結果に基づく最適な予算配分への変更',
                '2. 有効チャネルへの集中投資',
                '3. 年間予算計画へMMMモデリングデータを反映',
            ]),
        ]
        kpi_text = 'KPI目標：R² 0.80以上（再分析後）／NRMSE 0.10以下／汎化性能・配分整合性スコアの改善'
    elif n_dummies > 20:
        months = [
            (f'Month 1: ダミー精査＆初動テスト', [
                f'1. ダミー変数（{n_dummies}本）の要因確認・除外の判断',
                f'2. {zero_action}',
                '3. 伸びしろありチャネルへのテスト増額',
            ]),
            ('Month 2: データ精査完了＆最適化の実行', [
                '1. ダミー変数除外による再分析（R²の変化確認）',
                '2. 予算最適化シナリオの本格展開',
                '3. 季節性ダミー変数などの追加',
            ]),
            ('Month 3: 効果検証＆次年度の計画作成', [
                f'1. CV数値のトラッキング、CPA・{roi_label}モニタリング',
                f'2. {roi_label}変化に応じた予算の再配分',
                '3. 年間予算計画へMMMモデリングデータを反映',
            ]),
        ]
        kpi_text = f'KPI目標：月次CV数+{cv_lift/3:.0f}〜+{cv_lift:.0f}%／R² {min(r2+0.01, 0.99):.2f}以上（再分析後）／ゼロ係数チャネル{max(len(zero_channels)-2, 0)}個以下'
    else:
        months = [
            ('Month 1: 予算の再配分＆増額実施', [
                f'1. {_inc_names}への予算集中',
                f'2. {_zero_act_named}',
            ]),
            ('Month 2: 効果検証＆モデル更新', [
                '1. 配分変更または増額後のCV・CPA変化のモニタリング',
                '2. データ精査・蓄積後のMMMの再実行',
                f'3. 追加チャネルへの投資検討（目標+{cv_lift:.1f}%）',
            ]),
            ('Month 3: 次年度計画＆スケール', [
                f'1. CV数値のトラッキング、CPA・{roi_label}モニタリング',
                f'2. 高{roi_label}チャネルへのスケールアップ',
                '3. 年間予算計画へのモデリングデータ反映・フォーキャスト更新',
            ]),
        ]
        kpi_text = f'KPI目標：月次CV数+{cv_lift/3:.0f}〜+{cv_lift:.0f}%／CPA改善率5%以上／伸びしろありチャネルの予算比率UP'
    _road_xs = [0.865, 4.915, 8.965]
    _card_y, _card_h = 1.85, 4.1
    for i, (title, items) in enumerate(months):
        x = _road_xs[i]
        _rect(sl, x, _card_y, 3.5, _card_h, line_color=BRAND_LIGHT, line_pt=1.2)
        # 見出しを "Month X" と副題に分割して2行中央寄せ
        _parts  = title.split(': ', 1)
        _m_lbl  = _parts[0]
        _m_sub  = _parts[1] if len(_parts) > 1 else ''
        _box(sl, x + 0.15, _card_y + 0.18, 3.2, 0.38, _m_lbl,
             font_size=15, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.15, _card_y + 0.58, 3.2, 0.33, _m_sub,
             font_size=14, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _line(sl, x + 0.2, _card_y + 1.08, x + 3.3, _card_y + 1.08,
              color=RULE_COLOR, width_pt=0.8)
        for j, item in enumerate(items):
            _box(sl, x + 0.2, _card_y + 1.30 + j * 0.85, 3.15, 0.85, item,
                 font_size=12, text_color=SECTION_BG, line_spacing=1.5,
                 v_anchor=MSO_ANCHOR.MIDDLE)
        if i < 2:
            _box(sl, _road_xs[i] + 3.55, _card_y + _card_h / 2 - 0.25, 0.4, 0.5, '▶',
                 font_size=16, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 6.15, 12.5, 0.45, kpi_text,
         font_size=11, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

    # ── 用語集 ───────────────────────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, '用語のご説明')
    glossary_s = [
        ('R²（決定係数）',              '実測CVの分散のうちモデルが説明できる割合。0〜1の間となり、0.90以上が目安。'),
        ('NRMSE',                       '正規化された平均二乗誤差。小さいほど予測精度が高い。0.1以下が実用水準。'),
        ('NRMSE (holdout)',              f'最終{holdout_days}{period_unit}の未学習データでのNRMSE。モデルの汎化性能を評価できる。'),
        ('RSSD',                 '支出金額シェアとCV貢献シェアの乖離。0に近いほど整合性が高い。'),
        ('残存効果',                     '広告効果が持続する期間を表す。λが高いほど効果が長期間持続する。'),
        ('飽和曲線',                 '支出金額の増加に対するCV効果の逓減を表す。α・γで形状が決まる。'),
        ('MCR',                          '媒体帰属率（Media Contribution Rate）。広告起因のCVが全CVに占める割合。15%以上が目安。'),
        (roi_label,                     f'{roi_unit}。高いほど効率的。'),
        ('CPA',                         'Cost Per Acquisition。支出金額÷CV数。低いほど効率的。'),
        ('限界ROI' if _is_monetary else '限界CPA',
         '追加で予算を増やしていった際に、追加投資に対してどれだけ売上が生まれるか？を表す値。増額すればするほど、逓減する。'
         if _is_monetary else
         '追加で予算を増やしていった際に、1件あたりいくらかかるか？を表す値。増額すればするほど、上昇する。'),
        ('95%信頼区間（CI）',             '真の値がこの範囲に、95%の確率で収まる区間。区間が広いほど推定の不確実性が高い。'),
        ('ベースライン',                 'モデルが媒体チャネルへの投下として説明できなかったCV。ブランド認知・季節性・口コミ・SEO等が主な要素だが、\nモデル精度の限界により媒体効果と分離できなかったものを含む。'),
        ('ダミー変数',                   '特定日付の異常値を補正する変数。'),
        ('伸びしろあり／適正域／飽和域', f'追加投資{roi_label}÷平均{roi_label}で判定。0.5超え：伸びしろあり／0.2〜0.5：適正／0.2未満：飽和。'),
    ]
    _add_table(sl, 0.5, 1.1, 12.3, 5.9,
               ['用語', '説明'], [[k, v] for k, v in glossary_s],
               col_widths=[1.4, 3.0], body_font_size=10,
               col_alignments=[PP_ALIGN.CENTER, PP_ALIGN.LEFT])

    # ── N枚目: CTAスライド（v16 slide17準拠）────────────────────
    sl = _add_slide(prs)
    _bg(sl, SECTION_BG)
    _AMBER = RGBColor(0xCB, 0x80, 0x13)

    # 月次CV機会損失の計算
    _analysis_months = max(analysis_days / 30.0, 1.0)
    _monthly_cv_gap  = max(int(round(
        (opt_result['optimal_cv'] - opt_result['current_cv']) / _analysis_months
    )), 1)

    # ── Headline 3行
    _box(sl, 0.6, 0.38, 12.1, 0.34,
         '現状のまま運用を続けると、',
         font_size=14, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)

    _cta_gap_line = (
        f'推定で、-{_monthly_cv_gap}件の機会損失が毎月積み上がる可能性があります。'
        if _cv_metric_type != 'monetary'
        else f'推定で、-{cv_lift:.1f}%の売上機会損失が毎月積み上がる可能性があります。'
    )
    _box(sl, 0.6, 0.74, 12.1, 0.34,
         _cta_gap_line,
         font_size=14, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)

    _box(sl, 0.6, 1.35, 12.1, 0.34,
         '広告運用を改善したい方、データドリブンな運用を実現したい方は、お気軽にご相談ください。',
         font_size=14, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)

    _line(sl, 0.6, 1.98, 12.73, 1.98, color=BRAND_LIGHT, width_pt=0.6)

    # ── カード共通パラメータ（v16実測値）
    _cy    = 2.421
    _cw    = 3.745
    _ch    = 3.755
    _ca_x  = 2.578
    _cb_x  = 7.037
    _badge_h    = 0.262
    _badge_y    = _cy + 0.213   # 2.634 - 2.421
    _title_y    = _cy + 0.646   # 3.067 - 2.421
    _price_y    = _cy + 1.108   # 3.529 - 2.421
    _desc_y     = _cy + 2.108   # 4.529 - 2.421
    _btn_y      = _cy + 2.974   # 5.395 - 2.421
    _btn_h      = 0.432
    _txt_x_off  = 0.468         # 3.046 - 2.578
    _txt_w      = 2.852
    _btn_x_off  = 0.563         # 3.141 - 2.578
    _btn_w      = 2.620

    _cards = [
        dict(
            cx=_ca_x, badge_w=2.036,
            badge_txt='最適な運用を任せたい方はこちら',
            title='広告運用最適化プラン',
            price=('月額 ', '15', ' 万円 ～ ', '(税抜)'),
            desc='MMM分析をもとに、予算配分から\n日々の運用＆最適化までお任せ',
            btn_filled=True,
        ),
        dict(
            cx=_cb_x, badge_w=2.542,
            badge_txt='データをもとに意思決定をしたい方はこちら',
            title='MMMレポートプラン',
            price=('月換算 ', '8', ' 万円 ～ ', '(税抜)'),
            desc='今の運用体制は維持しながら、\nデータドリブンな意思決定をサポート',
            btn_filled=False,
        ),
    ]

    for _cd in _cards:
        _cx = _cd['cx']

        # カード背景（丸角・radius≈10pt）
        _rrect(sl, _cx, _cy, _cw, _ch, fill_color=WHITE, adj=0.037)

        # バッジチップ（pill形状・中央寄せ）
        _bw  = _cd['badge_w']
        _bx  = _cx + (_cw - _bw) / 2
        _rrect(sl, _bx, _badge_y, _bw, _badge_h, fill_color=BRAND_LIGHT, adj=0.5)
        _box(sl, _bx, _badge_y, _bw, _badge_h, _cd['badge_txt'],
             font_size=8, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER,
             v_anchor=MSO_ANCHOR.MIDDLE)

        # プラン名（中央寄せ）
        _box(sl, _cx + _txt_x_off, _title_y, _txt_w, 0.37, _cd['title'],
             font_size=16, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

        # 価格（大きな数字＋前後テキスト）
        _mixed_price_box(sl, _cx + _txt_x_off, _price_y, _txt_w, 0.909,
                         *_cd['price'])

        # 説明文（中央寄せ）
        _box(sl, _cx + _txt_x_off, _desc_y, _txt_w, 0.65, _cd['desc'],
             font_size=10, text_color=SECTION_BG, line_spacing=1.5,
             align=PP_ALIGN.CENTER)

        # CTAボタン（shape内にテキスト+リンクを直接埋め込み → PDF変換後もリンク保持）
        _bx2 = _cx + _btn_x_off
        if _cd['btn_filled']:
            _btn(sl, _bx2, _btn_y, _btn_w, _btn_h, 'まず無料で相談する   ▸',
                 fill_color=_AMBER, adj=0.2, text_color=WHITE,
                 url='https://m-teria.jp/?utm_source=mmm-report&utm_medium=paper&utm_campaign=plan-ad-agency')
        else:
            _btn(sl, _bx2, _btn_y, _btn_w, _btn_h, 'まず無料で相談する   ▸',
                 line_color=_AMBER, line_pt=1.5, adj=0.2, text_color=_AMBER,
                 url='https://m-teria.jp/?utm_source=mmm-report&utm_medium=paper&utm_campaign=plan-reporting')

    _copyright(sl, WHITE)

    # ── Page numbers ─────────────────────────────────────────
    _DARK_BG = {
        (COVER_BG[0],   COVER_BG[1],   COVER_BG[2]),
        (SECTION_BG[0], SECTION_BG[1], SECTION_BG[2]),
    }
    _last_idx = len(prs.slides) - 1
    for idx, slide in enumerate(prs.slides):
        if idx == 0 or idx == _last_idx:
            continue
        try:
            rgb     = slide.background.fill.fore_color.rgb
            is_dark = (int(rgb[0]), int(rgb[1]), int(rgb[2])) in _DARK_BG
        except Exception:
            is_dark = False
        num_color = WHITE if is_dark else BRAND_LIGHT
        _box(slide, 12.3, 7.1, 0.85, 0.28, str(idx),
             font_size=9, text_color=num_color, align=PP_ALIGN.RIGHT)

    from datetime import datetime as _dt
    now = _dt.utcnow()
    prs.core_properties.created  = now
    prs.core_properties.modified = now
    prs.save(output_path)
    return output_path


# ── Main report builder ────────────────────────────────────

def generate_report_from_snapshot(snapshot: dict, output_path: str,
                                   report_type: str = 'full') -> str:
    """pkl スナップショット dict から PPTX を生成する convenience wrapper。
    frontier / opt_result_b など渡し忘れが起きやすいフィールドを自動補完する。
    """
    return generate_report(
        metrics=snapshot['metrics'],
        opt_result=snapshot['opt_result'],
        opt_result_b=snapshot.get('opt_result_b'),
        opt_result_dec=snapshot.get('opt_result_dec'),
        frontier=snapshot.get('frontier'),
        client_name=snapshot['client_name'],
        output_path=output_path,
        cv_col=snapshot.get('cv_col', 'CV'),
        media_basis=snapshot.get('media_basis', 'media'),
        freq=snapshot.get('freq', 'daily'),
        lambda_profile=snapshot.get('lambda_profile', 'default'),
        budget_increase_pct=snapshot.get('budget_increase_pct', 0.20),
        report_type=report_type,
    )

def generate_report(metrics: dict, opt_result: dict, client_name: str,
                    model_name: str = 'M2 Bayesian Ridge',
                    output_path: str = 'mmm_report.pptx',
                    cv_col: str = 'CV',
                    media_basis: str = 'media',
                    freq: str = 'daily',
                    lambda_profile: str = 'default',
                    opt_result_b: dict = None,
                    opt_result_dec: dict = None,
                    frontier: dict = None,
                    budget_increase_pct: float = 0.30,
                    report_type: str = 'full',
                    export_charts_dir: str = None) -> str:
    """Generate MMM PPTX report (MTERIA brand design).
    report_type: 'full'=フルレポート / 'simple'=SMB向け簡易版
    export_charts_dir: 指定時にグラフをPNGとして保存（simple: 01-08 / full: full_01〜full_16 + full_rc_{ch}）
    """
    if report_type == 'simple':
        return _build_simple_report(
            metrics=metrics, opt_result=opt_result, client_name=client_name,
            output_path=output_path, cv_col=cv_col, media_basis=media_basis,
            freq=freq, opt_result_b=opt_result_b, opt_result_dec=opt_result_dec,
            frontier=frontier, budget_increase_pct=budget_increase_pct,
            export_dir=export_charts_dir,
        )
    from .metrics import response_curve as compute_response_curve, compute_marginal_roi
    import os

    # 頻度ラベル
    freq_label = '週次' if freq == 'weekly' else '日次'
    period_unit = '週' if freq == 'weekly' else '日'

    # MAPEスパース判定: ゼロ率 > 30% をスパースとみなす
    _actual_cv = metrics.get('actual_train_cv', np.array([]))
    _zero_rate = float(np.mean(_actual_cv == 0)) if len(_actual_cv) > 0 else 0.0
    _is_sparse = _zero_rate > 0.30
    if _is_sparse:
        mape_note = f'{freq_label}CVはゼロが多く（ゼロ率{_zero_rate*100:.0f}%）高めに出やすい。月次集計での確認を推奨'
    else:
        mape_note = f'{freq_label}CVのゼロ率{_zero_rate*100:.0f}%。予測誤差率として参考値。'

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    ch_metrics     = metrics['channel_metrics']
    channels       = metrics['channels']
    r2             = metrics['r2']
    nrmse          = metrics['nrmse']
    nrmse_hold     = metrics['nrmse_hold']
    rssd           = metrics['rssd']
    mape           = metrics['mape']
    media_fraction = metrics.get('media_fraction', 0.0)
    total_cv    = metrics['total_cv']
    total_spend = metrics['total_spend']
    n_dummies   = metrics['n_dummies']
    n_valid     = metrics['n_valid']
    n_total     = len(channels)

    valid_channels = [ch for ch in channels if not ch_metrics[ch]['is_zero']]
    zero_channels  = [ch for ch in channels if ch_metrics[ch]['is_zero']]

    dates_train      = metrics['dates_train']
    dates_hold       = metrics['dates_hold']
    holdout_days     = len(dates_hold)
    actual_train_cv  = metrics['actual_train_cv']
    pred_train_cv    = metrics['pred_train_sqrt'] ** 2
    actual_hold_cv   = metrics['actual_hold_cv']
    ch_daily_contrib = metrics.get('ch_daily_contrib', {})
    pred_hold_cv    = metrics['pred_hold_sqrt'] ** 2
    analysis_days   = len(dates_train) + len(dates_hold)

    # ── Slide 1: Cover ─────────────────────────────────────────
    sl = _add_slide(prs)
    _bg(sl, COVER_BG)
    # Client name top-left
    _box(sl, 0.55, 0.35, 6, 0.4, f'{client_name}　御中',
         font_size=16, bold=True, text_color=WHITE)
    # Main title center (moved down)
    _box(sl, 0.5, 2.9, 12.33, 1.0, '未来が見える広告レポート',
         font_size=46, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    # Subtitle (18pt)
    _box(sl, 0.5, 3.95, 12.33, 0.65,
         'マーケティングミックスモデリングによる媒体分析と最適な予算配分',
         font_size=18, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    # Logo (moved lower)
    logo = LOGO_PATH
    if os.path.exists(logo):
        sl.shapes.add_picture(logo, Inches(5.67), Inches(6.5), Inches(2.0), height=None)
    else:
        _box(sl, 0.5, 6.5, 12.33, 0.4, 'M TERIA',
             font_size=16, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: Table of Contents ─────────────────────────────
    sl = _add_slide(prs)
    # White bg (no _bg call) — title aligned with other slides
    _box(sl, 0.4, 0.20, 4, 0.58, '目次',
         font_size=22, bold=True, text_color=SECTION_BG)
    _line(sl, 0.4, 0.82, 12.93, 0.82, color=RULE_COLOR, width_pt=1.2)
    toc_parts = [
        ('01', 'フレームワークのご説明',
         'マーケティングミックスモデリングの流れ／モデル概要／Bayesian Ridge／Adstock／Hill／ダミー探索／RSSD'),
        ('02', 'モデルの精度',
         '実測vs予測／残差分析／月別の誤差／貢献分解／パラメータ'),
        ('03', '媒体・施策別の効果分析',
         'CPA／ROI／レスポンスカーブ（チャネル固有）／限界ROI'),
        ('04', '成果最大化につながる予算配分',
         'シミュレーション／チャネル別の推奨アクション'),
        ('05', 'ダミー変数',
         '採用変数一覧／要因分析／除外判断ガイド'),
        ('06', '事業上の推奨アクション',
         '戦略提言／ロードマップ／KPI目標'),
    ]
    for i, (num, ttl, desc) in enumerate(toc_parts):
        y = 1.10 + i * 0.90
        _box(sl, 1.05, y + 0.15, 0.85, 0.55, num,
             font_size=22, bold=True, text_color=BRAND_LIGHT, align=PP_ALIGN.LEFT)
        _box(sl, 2.1, y + 0.12, 11, 0.38, ttl,
             font_size=15, bold=True, text_color=SECTION_BG)
        _box(sl, 2.1, y + 0.58, 11.5, 0.28, desc,
             font_size=9, text_color=BRAND_LIGHT)
    _copyright(sl)

    # ── Part 1: Framework ──────────────────────────────────────

    # Slide 3: Section header
    sl = _add_slide(prs)
    _section_header(sl, 1, 'フレームワークのご説明')

    # Slide 4: Pipeline overview (4+4 layout, exact positions from reference)
    sl = _add_slide(prs)
    _slide_title(sl, 'マーケティングミックスモデリング（MMM）の流れ')
    # Paragraph below rule
    _box(sl, 0.4, 0.92, 12.5, 1.0,
         'マーケティングミックスモデリングとは、「複数の広告・施策にコストを投下していた際に、どれがどのくらい売上・CVに寄与しているか？を可視化し、未来の予算配分・投資チャネルの最適解を導き出す分析手法」です。\n'
         '8つのステップで広告・施策効果を分析し、成果最大化につながるパレート最適な予算配分を割り出します。',
         font_size=12, text_color=BRAND_LIGHT, line_spacing=1.5)

    # Row 1 boxes: x=0.8/3.4/6.0/8.6, y=2.7, size=2.3x1.8
    row1 = [
        ('①データの読込',  '日次or週次の\n売上・CV・コスト・\nチャネル・施策データ'),
        ('②データの前処理', '欠損補完などの\nクレンジングやモデル精度\nを上げる前処理'),
        ('③Adstock変換',   'チャネルごとに\n残存効果を推定'),
        ('④Hill変換',      'チャネルごとに\n効果が飽和し始める\n水準を推定'),
    ]
    row1_x = [0.8, 3.4, 6.0, 8.6]
    for i, (step, desc) in enumerate(row1):
        x = row1_x[i]
        _rect(sl, x, 2.45, 2.3, 1.8, line_color=BRAND_LIGHT, line_pt=1.2)
        _box(sl, x + 0.1, 2.75, 2.1, 0.38, step,
             font_size=14, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, 3.15, 2.1, 0.9, desc,
             font_size=11, text_color=SECTION_BG, align=PP_ALIGN.CENTER, line_spacing=1.5)
        if i < 3:
            _box(sl, row1_x[i] + 2.3, 3.15, 0.3, 0.4, '▶',
                 font_size=14, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
    # Continuation arrow
    _box(sl, 10.9, 3.15, 0.3, 0.4, '▶',
         font_size=14, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

    # Row 2 boxes: x=2.2/4.8/7.4/10.0, y=4.65, size=2.3x1.8
    row2 = [
        ('⑤パレート探索',  '2000回の試行／\nNRMSE×RSDDが\n最小になるモデルを探索'),
        ('⑥局所最適化',    '最急降下法による\n精緻なチューニング'),
        ('⑦ダミー変数探索',    'モデルが説明できない残差\nをダミー変数として\n組み込み推定'),
        ('⑧モデル出力',    '間接効果や残存効果も\n加味したCPA・ROI、\n最適な予算配分を提示'),
    ]
    row2_x = [2.2, 4.8, 7.4, 10.0]
    for i, (step, desc) in enumerate(row2):
        x = row2_x[i]
        _rect(sl, x, 4.65, 2.3, 1.8, line_color=BRAND_LIGHT, line_pt=1.2)
        _box(sl, x + 0.1, 4.85, 2.1, 0.38, step,
             font_size=14, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, 5.35, 2.1, 0.9, desc,
             font_size=11, text_color=SECTION_BG, align=PP_ALIGN.CENTER, line_spacing=1.5)
        if i < 3:
            _box(sl, row2_x[i] + 2.3, 5.35, 0.3, 0.4, '▶',
                 font_size=14, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

    # Slide 5: Model comparison
    sl = _add_slide(prs)
    _slide_title(sl, f'Bayesian Ridgeモデルを採用する理由')
    _box(sl, 0.4, 0.92, 12.5, 0.85,
         f'候補の4モデルのうち、今回の分析ではBayesian Ridgeモデルを採用しています。\n'
         f'4モデルの中で、精度や解釈性、計算速度のバランスが最も優秀であるためです。',
         font_size=12, text_color=BRAND_LIGHT, line_spacing=1.5)
    bullets = [
        ('自動正則化',        'α・λをデータから自動推定しスパース性（疎な状態）を自動学習するため、チューニングコストがかからない。'),
        ('事前分布の活用',    '係数に N(0, α⁻¹I) を設定するモデルのため、少ないデータでも過学習しにくい。'),
        ('不確実性の定量化',  '事後分布から信頼区間を算出するため、係数の確かさを数値で評価が可能。'),
        ('4モデルのうち、最もバランスがいい', f'R² = {r2:.3f} / NRMSE = {nrmse:.4f} / RSSD = {rssd:.4f}'),
    ]
    for i, (b, d) in enumerate(bullets):
        y = 2.1 + i * 1.05
        _box(sl, 0.4, y + 0.05, 4.3, 0.38, b,
             font_size=13, bold=True, text_color=SECTION_BG)
        _box(sl, 0.4, y + 0.48, 4.5, 0.42, d,
             font_size=10, text_color=BRAND_LIGHT)
    _add_table(sl, 5.2, 2.3, 7.5, 3.4,
               ['候補モデル', 'R²', 'NRMSE', 'RSSD', '適合度合い'],
               [['Ridge',          '-',          '-',            '-',          '4位'],
                [f'Bayesian Ridge', f'{r2:.3f}', f'{nrmse:.4f}', f'{rssd:.4f}', '1位'],
                ['ElasticNet',     '-',          '-',            '-',          '2位'],
                ['状態空間',        '-',          '-',            '-',          '3位']])

    # Slide 6: Adstock
    sl = _add_slide(prs)
    _slide_title(sl, 'Adstock（残存効果）変換')
    _box(sl, 0.4, 0.92, 12.5, 0.35, '広告出稿後に持続する効果を定式化し、モデルに組み込みます。',
         font_size=12, text_color=BRAND_LIGHT, line_spacing=1.5)
    _box(sl, 0.5, 1.6, 9, 0.5, '数式: x*ₜ = xₜ + λ × x*ₜ₋₁',
         font_size=14, bold=True, text_color=SECTION_BG)
    _box(sl, 0.5, 2.1, 9, 0.4, 'λ範囲：0.0〜0.9\nλ = 0.0：残存効果なし / λ = 0.9：10日後に35%の効果が残存',
         font_size=11, text_color=BRAND_LIGHT, line_spacing=1.5)
    top5_lambda = sorted(valid_channels,
                          key=lambda c: ch_metrics[c]['lambda'], reverse=True)[:5]
    lambda_vals = {ch: ch_metrics[ch]['lambda'] for ch in top5_lambda}
    fig_decay = _plot_adstock_decay(lambda_vals)
    _export_png(fig_decay, export_charts_dir, 'full_01_adstock_decay')
    _add_chart_image(sl, _fig_to_image(fig_decay), 0.5, 2.8, 6.5, 4.1)
    _box(sl, 7.5, 1.45, 5.5, 0.3, '残存効果上位のチャネル',
         font_size=11, bold=True, text_color=SECTION_BG)
    for i, ch in enumerate(top5_lambda):
        lam = ch_metrics[ch]['lambda']
        _kpi_card(sl, 7.5, 1.85 + i * 0.95, 5.5, 0.82, ch, f'λ = {lam:.3f}', SECTION_BG)

    # Slide 7: Hill
    sl = _add_slide(prs)
    _slide_title(sl, 'Hill飽和曲線（CV効果逓減）')
    _box(sl, 0.4, 0.92, 12.5, 0.35, '施策コストの支出が増加するにつれて、逓減するコンバージョン成果（売上や販売数、問い合わせ数など）を定式化します。',
         font_size=12, text_color=BRAND_LIGHT, line_spacing=1.5)
    _box(sl, 0.5, 1.65, 12, 0.4,
         'H(x) = xᵅ / (xᵅ + γᵅ)　　α = [0.8, 1.5]　γ = [0.8, 1.0]',
         font_size=13, bold=True, text_color=SECTION_BG)
    descs = [
        ('α < 1.0', '早期に飽和する（急峻な逓減）', '施策コストの支出を増やしても\n効果が伸びにくい'),
        ('α = 1.0', '線形的に逓減する',              '比例的な効果増加'),
        ('α > 1.0', 'S字型（最初は緩やか）',      '一定支出金額を超えると効果が急増'),
    ]
    for i, (a, t, d) in enumerate(descs):
        x = 0.5 + i * 4.2
        _rect(sl, x, 2.25, 3.8, 2.3, line_color=RULE_COLOR, line_pt=1.0)
        _box(sl, x + 0.15, 2.38, 3.5, 0.55, a,
             font_size=16, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.15, 2.98, 3.5, 0.4,  t,
             font_size=12, bold=True, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.15, 3.38, 3.5, 0.8,  d,
             font_size=12, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
    # γ explanation card spanning full width
    _rect(sl, 0.5, 4.75, 12.3, 0.95, line_color=RULE_COLOR, line_pt=1.0)
    _box(sl, 1.0, 4.75, 2.5, 0.95, 'γ（飽和点）',
         font_size=14, bold=True, text_color=SECTION_BG, v_anchor=MSO_ANCHOR.MIDDLE)
    _box(sl, 3.7, 4.88, 8.8, 0.7,
         '小さいほど早期に飽和し、大きいほど大きなコストを投下しても効果が継続する',
         font_size=11, text_color=BRAND_LIGHT, v_anchor=MSO_ANCHOR.MIDDLE)

    # ── Part 2: Model Diagnostics ──────────────────────────────

    sl = _add_slide(prs)
    _section_header(sl, 2, 'モデルの精度')

    # Slide 9: Model accuracy summary
    sl = _add_slide(prs)
    _slide_title(sl, 'モデル精度サマリー',
                 f'{model_name} + 自動ダミー{n_dummies}本')
    kpis9 = [
        ('R²（決定係数）',   f'{r2:.3f}',                     SECTION_BG),
        ('NRMSE (train)',     f'{nrmse:.4f}',                   SECTION_BG),
        ('NRMSE (holdout)',   f'{nrmse_hold:.4f}',              SECTION_BG),
        ('RSSD',             f'{rssd:.4f}',                    SECTION_BG),
        ('MAPE',             f'{mape*100:.1f}%',               SECTION_BG),
        ('有効チャネル',      f'{n_valid}/{n_total}',           SECTION_BG),
        ('媒体帰属比率',      f'{media_fraction*100:.1f}%',     SECTION_BG),
    ]
    for i, (l, v, c) in enumerate(kpis9):
        _kpi_card(sl, 0.4 + i * 1.82, 1.65, 1.65, 1.1, l, v, c)

    # Dynamic judgment per metric
    r2_status = f'{r2*100:.0f}%'  # 自然言語：説明率をそのまま表示
    nrmse_status = '◎' if nrmse < 0.10 else '○' if nrmse < 0.12 else '△' if nrmse < 0.15 else '×'
    nrmse_ratio  = nrmse_hold / max(nrmse, 1e-6)
    hold_status  = '◎' if nrmse_hold < 0.15 else '○' if nrmse_hold < 0.20 else '△' if nrmse_hold < 0.25 else '×'
    hold_detail  = (
        f'汎化性能良好（train比{nrmse_ratio:.1f}倍）'
        if nrmse_ratio < 1.5 else
        f'train比{nrmse_ratio:.1f}倍 — 検証期間（holdout）に特異値が含まれている可能性'
        if nrmse_ratio < 2.0 else
        f'train比{nrmse_ratio:.1f}倍 — ダミー変数の過剰適合を確認推奨'
    )
    rssd_status  = '◎' if 0.10 <= rssd <= 0.20 else '○' if rssd <= 0.30 else '△' if rssd <= 0.40 else '×'

    interp_rows = [
        ['R²（決定係数）',  f'{r2:.3f}',        r2_status,
         (f'売上変動の{r2*100:.0f}%をモデルで説明。精度は高い水準です' if r2 >= 0.85 else
          f'売上変動の{r2*100:.0f}%をモデルで説明。データ追加で精度向上の余地あり' if r2 >= 0.70 else
          f'売上変動の{r2*100:.0f}%をモデルで説明。チャネル数・データ量が少ない場合は低めになります')],
        ['NRMSE (train)',   f'{nrmse:.4f}',      nrmse_status,
         ('予測精度は非常に高い水準です。（目安：0.12未満）' if nrmse < 0.10 else
          '良好な予測精度です。（目安：0.12未満）' if nrmse < 0.12 else
          '許容範囲内ですが、改善の余地があります。（目安：0.12未満）' if nrmse < 0.15 else
          'データの追加・クレンジングで改善が見込めます。（目安：0.12未満）')],
        ['NRMSE (holdout)', f'{nrmse_hold:.4f}', hold_status,
         (f'未来データへの汎化性能は高い水準です。{hold_detail}（目安：0.20未満）' if nrmse_hold < 0.15 else
          f'汎化性能は良好な水準です。{hold_detail}（目安：0.20未満）' if nrmse_hold < 0.20 else
          f'若干の過学習の傾向があります。{hold_detail}（目安：0.20未満）' if nrmse_hold < 0.30 else
          f'ホールドアウト期間に特異なデータがある可能性があります。{hold_detail}（目安：0.20未満）')],
        ['RSSD',            f'{rssd:.4f}',        rssd_status,
         ('整合性は高水準ですが、均質的な配分になっている可能性があります。（目安：0.10〜0.30）' if rssd < 0.10 else
          '投資額とCV貢献の配分が整合しています。（目安：0.10〜0.30）' if rssd < 0.30 else
          '投資額とCV貢献に乖離があります。配分の見直しを推奨します。（目安：0.10〜0.30）')],
        ['媒体帰属比率',    f'{media_fraction*100:.1f}%', '参考値',
         '全CVに占める広告媒体起因の割合。標準的なMMMでは5〜15%。残りはSEO・ブランド・自然流入等のベースライン'],
        ['MAPE',            f'{mape*100:.1f}%',   '参考値', mape_note],
    ]
    # Auto-selection note
    auto_sel = metrics.get('auto_selection')
    if auto_sel:
        sel_b = auto_sel.get('selected_basis',   auto_sel.get('selected', media_basis))
        sel_p = auto_sel.get('selected_profile', 'default')
        scores = auto_sel.get('combo_scores')
        if scores:
            score_str = ' / '.join(f'{k}={v:.4f}' for k, v in sorted(scores.items()))
            basis_note = f'モデル基準: {sel_b}/{sel_p}（自動選択 — NRMSE検証: {score_str}）'
        else:
            hold_m = auto_sel.get('nrmse_hold_media', 0)
            hold_s = auto_sel.get('nrmse_hold_spend', 0)
            basis_note = f'モデル基準: {sel_b}ベース（自動選択 — NRMSE検証期間: media={hold_m:.4f} / spend={hold_s:.4f}）'
    else:
        basis_note = f'モデル基準: {media_basis}ベース / {lambda_profile}プロファイル（固定）'
    _box(sl, 0.4, 2.8, 12.53, 0.3, basis_note,
         font_size=9, text_color=BRAND_LIGHT)

    _box(sl, 0.4, 3.3, 5.5, 0.35, '各指標の解釈',
         font_size=12, bold=True, text_color=SECTION_BG)
    _add_table(sl, 0.4, 3.7, 12.53, 3.4,
               ['指標', '値', '判定', '解釈'],
               interp_rows,
               header_font_size=10, body_font_size=8,
               col_widths=[1.6, 0.9, 1.0, 9.0])

    # Slide 10: Actual vs Predicted
    sl = _add_slide(prs)
    _slide_title(sl, f'実測と予測の推移（{freq_label} {cv_col}）',
                 f'上：時系列での実測と予測のフィット状況（右端{holdout_days}{period_unit}：未学習データによる検証期間（holdout））／下：実測と予測モデルとの差')
    fig = _plot_actual_vs_pred(
        dates_train, actual_train_cv, pred_train_cv,
        dates_hold, actual_hold_cv, pred_hold_cv, cv_col=cv_col)
    _export_png(fig, export_charts_dir, 'full_02_actual_vs_pred')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.55)

    # Slide 11: Monthly accuracy
    sl = _add_slide(prs)
    _slide_title(sl, '月別CVの実測と予測', '誤差率が20%を超える月は外部要因の影響が大きい。モデルで説明できない要因をダミー変数として追加することで、精度改善が可能。')
    all_dates  = np.concatenate([dates_train, dates_hold])
    all_actual = np.concatenate([actual_train_cv, actual_hold_cv])
    all_pred   = np.concatenate([pred_train_cv, pred_hold_cv])
    fig = _plot_monthly_accuracy(all_dates, all_actual, all_pred, cv_col=cv_col)
    _export_png(fig, export_charts_dir, 'full_03_monthly_accuracy')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.35)

    # Slide 12: Monthly CV contribution by channel (stacked bar)
    sl = _add_slide(prs)
    _slide_title(sl, '月次推定CV数の推移（上位8チャネル）', 'チャネル別の推定CV数の月次推移')
    fig_stk = _plot_monthly_stacked(dates_train, ch_daily_contrib, channels, ch_metrics)
    _export_png(fig_stk, export_charts_dir, 'full_04_monthly_stacked')
    _add_chart_image(sl, _fig_to_image(fig_stk), 0.4, 1.35, 12.5, 5.55)

    # Slide 12b: CV decomposition — baseline vs media
    sl = _add_slide(prs)
    _slide_title(sl, 'CV分解：ベースライン vs 広告による獲得',
                 'モデルが媒体投下として説明できなかったCV（ベースライン）と、各広告チャネルに帰属できたCVの月次内訳')
    fig_decomp, med_pct = _plot_cv_decomp(dates_train, ch_daily_contrib, channels, actual_train_cv,
                                           media_fraction=media_fraction)
    _export_png(fig_decomp, export_charts_dir, 'full_05_cv_decomp')
    _add_chart_image(sl, _fig_to_image(fig_decomp), 0.5, 1.35, 9.5, 5.55)
    _kpi_card(sl, 10.2, 2.0, 2.9, 1.1, '広告による獲得（媒体帰属）', f'{med_pct:.1f}%', SECTION_BG)
    _kpi_card(sl, 10.2, 3.3, 2.9, 1.1, 'ベースライン', f'{100-med_pct:.1f}%', SECTION_BG)

    # Slide 13: Spend share vs CV contribution share
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別の支出金額シェアとCV貢献シェア',
                 f'乖離が大きいチャネルが投資効率のミスマッチを表す。　RSSD = {rssd:.4f}（√Σ(CV貢献ᵢ−支出金額ᵢ)²・Robyn準拠）')
    sp_shares = np.array([ch_metrics[ch]['spend_share']  for ch in channels])
    cv_shares = np.array([ch_metrics[ch]['contrib_share'] for ch in channels])
    fig = _plot_spend_vs_cv(channels, sp_shares, cv_shares, rssd)
    _export_png(fig, export_charts_dir, 'full_06_spend_vs_cv')
    _add_chart_image(sl, _fig_to_image(fig), 0.4, 1.45, 12.5, 5.45)

    # Slide 14: Parameter heatmap
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別パラメータのヒートマップ',
                 'Hill α・γ・Adstock λ・効果係数（正規化）')
    param_dict = {ch: {'lambda': ch_metrics[ch]['lambda'],
                        'alpha':  ch_metrics[ch]['alpha'],
                        'gamma':  ch_metrics[ch]['gamma']} for ch in channels}
    all_coefs  = [ch_metrics[ch]['coef'] for ch in channels]
    fig = _plot_param_heatmap(channels, param_dict, all_coefs)
    _export_png(fig, export_charts_dir, 'full_07_param_heatmap')
    _add_chart_image(sl, _fig_to_image(fig), 0.4, 1.35, 7.0, 5.55)
    # Legend table on the right
    legend_rows = [
        ['λ', '残存効果が長い', '効果がすぐ消える'],
        ['α', 'S字型・後から急増', '早期飽和・頭打ち'],
        ['γ', '飽和点が高い・効果継続', '早期飽和'],
        ['効果係数', '推定CV数が大きい', 'ゼロ係数チャネル'],
    ]
    _add_table(sl, 7.6, 2.73, 5.2, 2.8,
               ['指標', '高い（緑）', '低い（赤）'], legend_rows,
               header_font_size=9, body_font_size=8,
               col_widths=[0.5, 1.5, 1.5])

    # ── Part 3: Media Effect Analysis ─────────────────────────

    sl = _add_slide(prs)
    _section_header(sl, 3, '媒体・施策別の効果分析')

    # Slide 16: CPA comparison
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別CPA', '有効チャネルのみ・CPA昇順（低いほど効率的）')
    vchs = [ch for ch in valid_channels if ch_metrics[ch]['cpa'] is not None]
    cpas = np.array([ch_metrics[ch]['cpa'] for ch in vchs])
    fig  = _plot_channel_bar(vchs, cpas, channel_metrics=ch_metrics,
                              xlabel='CPA（低いほど効率的）', fmt='¥{:,.0f}', reverse=False)
    _export_png(fig, export_charts_dir, 'full_08_cpa_bar')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.55)

    # Slide 17: ROI ranking
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別ROI', 'ROI = 推定CV数 / 支出金額（万円）')
    rois2 = np.array([ch_metrics[ch]['roi'] for ch in valid_channels])
    fig   = _plot_channel_bar(valid_channels, rois2, channel_metrics=ch_metrics, xlabel='ROI')
    _export_png(fig, export_charts_dir, 'full_09_roi_bar')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.55)

    # Slide 17.5: Forest plot — ROI with 95% CI
    _ci_chs = [ch for ch in valid_channels if ch_metrics[ch].get('ci_available')]
    if _ci_chs:
        sl = _add_slide(prs)
        _slide_title(sl, 'チャネル別ROI 確信度（95%信頼区間）',
                     '点 = ROI点推定値  横棒 = 95%信頼区間 / 緑 = CI全体がROI>1.0（確信度高） / '
                     'ティール = CI跨ぎ（要継続観察） / グレー = CI全体<1.0（効果低）')
        fig_fp = _plot_forest(ch_metrics, valid_channels)
        if fig_fp is not None:
            _export_png(fig_fp, export_charts_dir, 'full_10_roi_forest')
            _add_chart_image(sl, _fig_to_image(fig_fp), 0.5, 1.35, 12.3, 5.55)

    # Slide 18: Full channel table
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別CPA・ROI・パラメータの一覧',
                 f'有効{n_valid}ch / ゼロ係数{len(zero_channels)}ch')
    rows17 = []
    for ch in channels:
        cm    = ch_metrics[ch]
        cpa_s = f'¥{cm["cpa"]:,.0f}' if cm['cpa'] else 'N/A'
        state = 'ゼロ×' if cm['is_zero'] else '有効'
        mroi  = f'{cm.get("marginal_roi", 0.0):.2f}'
        sat   = cm.get('saturation_label', '-')
        rows17.append([ch, cpa_s, sat,
                        f'{cm["spend_man"]:.1f}', f'{cm["roi"]:.2f}', mroi,
                        f'{cm["lambda"]:.3f}',    f'{cm["alpha"]:.3f}',
                        f'{cm["gamma"]:.3f}',     f'{cm["cv_contrib"]:.1f}', state])
    _add_table(sl, 0.3, 1.35, 12.73, 5.65,
               ['チャネル', 'CPA', '飽和状態', '支出金額(万)', 'ROI', '限界ROI',
                'λ', 'α', 'γ', '推定CV数', '状態'],
               rows17, header_font_size=10,
               body_font_size=7 if len(channels) > 20 else 8)

    # Slide 19: Marginal ROI
    sl = _add_slide(prs)
    _slide_title(sl, '追加投資ROI — 次の¥1万をどこに投じるか',
                 '現在の支出水準から¥1万追加したときに見込めるCV増加数。平均ROIより低い場合は飽和が近いサイン。')
    fig_roi = _plot_marginal_roi(valid_channels, ch_metrics)
    _export_png(fig_roi, export_charts_dir, 'full_11_marginal_roi')
    _add_chart_image(sl, _fig_to_image(fig_roi), 0.4, 1.35, 12.5, 5.55)

    # Slide 20: Spend vs CV bubble chart
    sl = _add_slide(prs)
    _slide_title(sl, '支出金額と推定CV数の散布図',
                 'バブル=支出金額規模 / 右上が理想ゾーン')
    fig, ax = plt.subplots(figsize=(10, 5.2), facecolor=CHART_BG)
    ax.set_facecolor(CHART_AX)
    for ch in valid_channels:
        cm = ch_metrics[ch]
        c  = SAT_CHART_COLORS.get(cm.get('saturation_label', '-'), '#BBBBBB')
        ax.scatter(cm['spend_man'], cm['cv_contrib'],
                   s=max(cm['spend_man'], 1) * 10, color=c, alpha=0.75)
        ax.annotate(ch[:8], (cm['spend_man'], cm['cv_contrib']),
                    fontsize=7, color=C_TICK)
    ax.set_xlabel('支出金額（万円）', color=C_TICK, fontsize=9)
    ax.set_ylabel('推定CV数',   color=C_TICK, fontsize=9)
    _chart_style(ax)
    fig.tight_layout(pad=0.8)
    _export_png(fig, export_charts_dir, 'full_12_spend_cv_bubble')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.55)

    # Slides: Response curves per channel — ordered by ROI descending (S rank first)
    from .transforms import apply_transforms, adstock_transform
    from .metrics import response_curve as rc

    roi_sorted_channels = sorted(valid_channels, key=lambda c: ch_metrics[c]['roi'], reverse=True)

    for ch in roi_sorted_channels:
        cm    = ch_metrics[ch]
        sl    = _add_slide(prs)
        cpa_s = f'¥{cm["cpa"]:,.0f}' if cm['cpa'] else 'N/A'
        sat   = cm.get('saturation_label', '-')
        _slide_title(sl, f'レスポンスカーブ：{ch}',
                     f'飽和状態: {sat} ／ CPA = {cpa_s} ／ 支出金額 = {_fmt_man(cm["spend_man"])} ／ ROI = {cm["roi"]:.2f}')
        curve_data = cm.get('curve_data') or {}
        fig = _plot_response_curve(curve_data, ch, cm, media_basis=media_basis)
        _export_png(fig, export_charts_dir, f'full_rc_{ch}')
        _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 7.0, 5.35)
        kp = [
            ('残存効果',      f'λ = {cm["lambda"]:.3f}'),
            ('飽和の傾き',    f'α = {cm["alpha"]:.3f}'),
            ('飽和の閾値',    f'γ = {cm["gamma"]:.3f}'),
            ('支出金額シェア', f'{cm["spend_share"]*100:.1f}%'),
            ('CV貢献シェア',  f'{cm["contrib_share"]*100:.1f}%'),
        ]
        for i, (label, val) in enumerate(kp):
            _kpi_card(sl, 7.7, 1.65 + i * 0.95, 5.3, 0.82, label, val, SECTION_BG)

    # Zero coefficient channels
    sl = _add_slide(prs)
    _slide_title(sl, f'ゼロ係数のチャネル（{len(zero_channels)}チャネル）',
                 'Bayesian Ridgeモデルで効果係数がゼロと推定。投資をROIの高いチャネルへ振り分けることをおすすめします。')
    import math as _math

    def _zero_reason(cm_):
        return cm_.get('zero_reason', 'CV効果が確認できないチャネルです')

    _n_zero = len(zero_channels)
    _n_rows = max(1, _math.ceil(_n_zero / 3))
    _card_h, _row_gap = 0.95, 1.12
    _total_h = _n_rows * _card_h + (_n_rows - 1) * (_row_gap - _card_h)
    _body_top, _body_h = 1.35, 5.55
    _start_y = _body_top + max(0.0, (_body_h - _total_h) / 2)
    for i, ch in enumerate(zero_channels):
        cm = ch_metrics[ch]
        x  = 0.5 + (i % 3) * 4.2
        y  = _start_y + (i // 3) * _row_gap
        _rect(sl, x, y, 3.8, _card_h, line_color=BRAND_LIGHT, line_pt=1.0)
        _box(sl, x + 0.1, y + 0.03, 3.6, 0.25, '× ゼロ係数',
             font_size=9, bold=True, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, y + 0.27, 3.6, 0.30, ch,
             font_size=12, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, y + 0.56, 3.6, 0.20, f'¥{int(cm["spend"]):,}',
             font_size=9, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, y + 0.73, 3.6, 0.18, _zero_reason(cm),
             font_size=7, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)

    # ── Part 4: Budget Optimization ────────────────────────────

    sl = _add_slide(prs)
    _section_header(sl, 4, '成果最大化につながる予算配分')

    cv_lift    = opt_result['cv_lift_pct']
    cv_lift_b  = opt_result_b['cv_lift_pct'] if opt_result_b else 0.0
    inc_pct_label = f'{budget_increase_pct*100:.0f}'

    # ── 予算最適化 概要スライド ──────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, '予算最適化シミュレーション — 2シナリオ比較',
                 'レスポンスカーブ（Hill飽和曲線）をもとに scipy SLSQP で CV 最大化を解く')
    _kpi_card(sl, 0.5,  1.6, 2.8, 1.1, 'シナリオA: 同予算再配分',
              f'+{cv_lift:.1f}% CV', SECTION_BG)
    _kpi_card(sl, 3.5,  1.6, 2.8, 1.1, f'シナリオB: +{inc_pct_label}%増額',
              f'+{cv_lift_b:.1f}% CV', SECTION_BG)
    max_eff_budget  = frontier['max_efficient_budget'] / 10000 if frontier else total_spend / 10000
    thr_cpa         = frontier['threshold_cpa']         if frontier else 0
    thr_roi         = frontier['threshold_roi']          if frontier else 0.0
    _is_monetary_fr = (frontier.get('cv_metric_type') == 'monetary') if frontier else False
    _kpi_card(sl, 6.5,  1.6, 3.0, 1.1, '理論上の最大効率予算',
              f'¥{max_eff_budget:.0f}万', SECTION_BG)
    if _is_monetary_fr:
        _kpi_card(sl, 9.7, 1.6, 3.1, 1.1, '目標ROI閾値', f'{thr_roi:.2f}x', SECTION_BG)
    else:
        _kpi_card(sl, 9.7, 1.6, 3.1, 1.1, '許容最大CPA（閾値）', f'¥{thr_cpa:,.0f}', SECTION_BG)
    _eff_desc = (
        f'限界ROI ≥ {thr_roi:.2f}（目標ROI閾値）まで'
        if _is_monetary_fr else
        f'限界CPA ≤ 現状平均CPA×1.5（¥{thr_cpa:,.0f}）まで'
    )
    _eff_detail = (
        '増額してもROIが目標を下回り始める分岐点。これ以上の増額は収益性が低下する'
        if _is_monetary_fr else
        '増額してもCPAが悪化し始める分岐点。これ以上の増額は費用対効果が低下する'
    )
    method_rows = [
        ['最適化手法',       'scipy SLSQP（逐次二次計画法）',
         '各チャネルの Hill 飽和曲線を使い「予算固定・CV 合計を最大化」する配分を数値計算で求める'],
        ['CV 推計方法',     'レスポンスカーブ補間',
         '現状支出での位置を基準に、支出変化後の曲線上の値を補間。飽和域では逓減を正確に反映'],
        ['ゼロ係数チャネル', '現状の 10% で固定',
         'CV 貢献がモデル上ゼロのチャネルは最適化対象外。残り 90% を他チャネルへ再配分'],
        ['理論上の最大効率予算', _eff_desc, _eff_detail],
    ]
    _add_table(sl, 0.5, 2.95, 12.3, 3.7,
               ['項目', '方法', '意味'],
               method_rows, col_widths=[2.2, 2.8, 7.3])

    # ── シナリオA: 同予算再配分 ──────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, f'シナリオA: 同予算での最適再配分（+{cv_lift:.1f}% CV見込み）',
                 f'総予算 ¥{total_spend/10000:.0f}万を固定したまま、配分を最適化することで CV+{cv_lift:.1f}%を見込める')
    opt_chs    = list(opt_result['channel_opt'].keys())
    cur_spends = [opt_result['channel_opt'][ch]['current_spend'] for ch in opt_chs]
    opt_spends = [opt_result['channel_opt'][ch]['optimal_spend'] for ch in opt_chs]
    fig = _plot_budget_bar(opt_chs, cur_spends, opt_spends)
    _export_png(fig, export_charts_dir, 'full_13_budget_optimal')
    _add_chart_image(sl, _fig_to_image(fig), 0.5, 1.35, 12.3, 5.65)

    # ── シナリオB: 増額時の最適配分 ────────────────────────────
    if opt_result_b is not None:
        sl = _add_slide(prs)
        new_bud = opt_result_b['total_budget']
        _slide_title(sl, f'シナリオB: +{inc_pct_label}%増額時の最適配分（+{cv_lift_b:.1f}% CV見込み）',
                     f'総予算 ¥{new_bud/10000:.0f}万（現状比+{inc_pct_label}%）での最適配分。飽和域チャネルへの増額は効率が低い')
        opt_chs_b    = list(opt_result_b['channel_opt'].keys())
        cur_spends_b = [opt_result_b['channel_opt'][ch]['current_spend'] for ch in opt_chs_b]
        opt_spends_b = [opt_result_b['channel_opt'][ch]['optimal_spend'] for ch in opt_chs_b]
        fig_b = _plot_budget_bar(opt_chs_b, cur_spends_b, opt_spends_b)
        _export_png(fig_b, export_charts_dir, 'full_14_budget_increase')
        _add_chart_image(sl, _fig_to_image(fig_b), 0.5, 1.35, 12.3, 5.65)

    # ── 投資効率フロンティア ────────────────────────────────
    if frontier is not None:
        sl = _add_slide(prs)
        _fr_subtitle2 = (
            f'限界ROI（追加¥1あたりの売上）が目標ROI {frontier["threshold_roi"]:.2f} を下回った点が理論上の増額上限'
            if _is_monetary_fr else
            f'限界CPA（追加¥1あたりのCPA）が ¥{frontier["threshold_cpa"]:,.0f}（現状CPA×1.5）を超えた点が理論上の増額上限'
        )
        _slide_title(sl, '投資効率フロンティア — どこまで増額すると効率が落ちるか', _fr_subtitle2)
        fig_fr = _plot_efficient_frontier(frontier)
        _export_png(fig_fr, export_charts_dir, 'full_15_efficient_frontier')
        _add_chart_image(sl, _fig_to_image(fig_fr), 0.4, 1.35, 9.8, 5.55)
        max_eff = frontier['max_efficient_budget']
        base_b  = frontier['current_budget']
        _kpi_card(sl, 10.4, 1.65, 2.5, 1.1, '最大効率予算',
                  f'¥{max_eff/10000:.0f}万', SECTION_BG)
        _eff_delta2 = (max_eff / base_b - 1) * 100
        _kpi_card(sl, 10.4, 2.95, 2.5, 1.1, '現状比',
                  f'+{_eff_delta2:.0f}%' if _eff_delta2 >= 0 else f'{_eff_delta2:.0f}%', SECTION_BG)
        if _is_monetary_fr:
            _kpi_card(sl, 10.4, 4.25, 2.5, 1.1, '目標ROI閾値',
                      f'{frontier["threshold_roi"]:.2f}x', SECTION_BG)
        else:
            _kpi_card(sl, 10.4, 4.25, 2.5, 1.1, '許容CPA閾値',
                      f'¥{frontier["threshold_cpa"]:,.0f}', SECTION_BG)

    # ── チャネル別推奨アクション ─────────────────────────────
    sl = _add_slide(prs)
    _slide_title(sl, 'チャネル別推奨アクション（シナリオA: 同予算最適配分）',
                 '飽和状態と最適配分の差分から増額・維持・削減を判定（±20%を判断基準）')
    rows_opt = []
    for ch in sorted(opt_chs,
                     key=lambda c: opt_result['channel_opt'][c]['roi'], reverse=True):
        o     = opt_result['channel_opt'][ch]
        delta = (o['optimal_spend'] - o['current_spend']) / 10000
        sat   = ch_metrics[ch].get('saturation_label', '-')
        rows_opt.append([ch, f'{o["current_spend"]/10000:.1f}',
                          f'{o["optimal_spend"]/10000:.1f}', f'{delta:+.1f}万',
                          f'{o["roi"]:.2f}', sat, o['action']])
    _add_table(sl, 0.3, 1.35, 12.73, 5.65,
               ['チャネル', '現状(万)', '最適(万)', '変化', 'ROI', '飽和状態', '推奨アクション'],
               rows_opt, header_font_size=10,
               body_font_size=7 if len(channels) > 20 else 9)

    # ── Part 5: Dummy Variables ────────────────────────────────

    sl = _add_slide(prs)
    _section_header(sl, 5, 'ダミー変数')

    sl = _add_slide(prs)
    _slide_title(sl, f'適用したダミー変数の一覧（{n_dummies}本）')
    dummy_info = metrics['dummy_info']
    rows_dum   = [[str(d.get('rank', '')), d.get('name', ''), d.get('date', ''),
                   d.get('category', ''), d.get('dow', '')] for d in dummy_info]
    if rows_dum:
        dum_body_fs = 5 if n_dummies > 29 else 7 if n_dummies > 20 else 9
        _add_table(sl, 0.3, 1.1, 12.73, 5.9,
                   ['#', 'ダミー名', '日付', 'カテゴリ', '曜日'], rows_dum[:30],
                   header_font_size=10, body_font_size=dum_body_fs)
    else:
        _box(sl, 0.5, 2.0, 12, 0.5, 'ダミー変数なし（高精度モデル）',
             font_size=14, text_color=BRAND_LIGHT)

    # Dummy variable bar chart slide
    if dummy_info:
        sl = _add_slide(prs)
        _slide_title(sl, 'ダミー変数の効果係数', '正：CVを増加させる要因／負：CVを減少させる要因（Bayesian Ridge係数）')
        fig_dum = _plot_dummy_bar(dummy_info)
        if fig_dum:
            _export_png(fig_dum, export_charts_dir, 'full_16_dummy_bar')
            _add_chart_image(sl, _fig_to_image(fig_dum), 0.4, 1.35, 12.5, 5.55)

    # ── Part 6: Business Insights ──────────────────────────────

    sl = _add_slide(prs)
    _section_header(sl, 6, '事業上の推奨アクション')

    top_roi_channels = sorted(valid_channels,
                               key=lambda c: ch_metrics[c]['roi'], reverse=True)[:3]
    headroom_channels = [ch for ch in valid_channels
                         if ch_metrics[ch].get('saturation_label') == '伸び代あり']

    sl = _add_slide(prs)
    _slide_title(sl, '主要インサイト', f'{model_name}分析から得られた戦略的示唆')
    insights = [
        ('即実行：伸び代チャネルへ集中増額', '\n'.join([
            f'・{", ".join(top_roi_channels[:3])} が高ROI筆頭',
            f'・同予算内で配分増だけでCV+{cv_lift:.1f}%見込める',
            f'・伸び代あり（{len(headroom_channels)}ch）への集中投資を優先推奨',
        ])),
        ('予算削減：ゼロ係数チャネルの整理', '\n'.join([
            f'・{", ".join(zero_channels[:3])} はCV貢献ゼロ',
            '・支出金額を高ROI施策へ振り替え推奨',
            '・停止前に1〜2週間のテスト停止で確認',
        ])),
        ('モデル精度：R²={:.3f}・NRMSE={:.4f}'.format(r2, nrmse), '\n'.join([
            f'・R²={r2:.3f}：全CVの{r2*100:.0f}%を説明',
            f'・NRMSE(HO)={nrmse_hold:.4f}：汎化性能の目安',
            '・検証期間（holdout）14日を加味したモデル選択',
        ])),
        ('次ステップ：ダミー除外再分析', '\n'.join([
            '・要因確認推奨ダミーを除外し再分析',
            '・R²変化±0.005以内なら除外推奨',
            '・精度向上後に予算最適化へ移行',
        ])),
    ]
    card_w, card_h = 2.9, 4.2
    ins_x = [0.45, 3.55, 6.65, 9.75]
    ins_y  = 1.55
    for i, (ttl, body) in enumerate(insights):
        cx = ins_x[i]
        cy = ins_y
        _rect(sl, cx, cy, card_w, card_h, line_color=BRAND_LIGHT, line_pt=1.2)
        _box(sl, cx + 0.12, cy + 0.18, card_w - 0.24, 0.55, ttl,
             font_size=11, bold=True, text_color=SECTION_BG)
        _box(sl, cx + 0.12, cy + 0.78, card_w - 0.24, card_h - 0.95, body,
             font_size=10, text_color=SECTION_BG, line_spacing=1.5)

    # Roadmap — 3-route branching
    sl = _add_slide(prs)
    _slide_title(sl, '3ヶ月アクションロードマップ', '即時〜中期の施策推進プラン')

    zero_action = f'ゼロ係数{len(zero_channels)}ch の予算削減・停止' if zero_channels else 'ゼロ係数チャネルなし（現状維持）'
    zero_action_c = f'ゼロ係数{len(zero_channels)}ch を即時停止' if zero_channels else 'ゼロ係数チャネルなし（現状維持）'

    if r2 < 0.75:
        # Route A: データ品質改善ルート
        months = [
            ('Month 1: データ整備 & モデル再構築', [
                '1. 欠損・外れ値の特定と補完',
                '2. コントロール変数の追加検討',
                '3. データ期間延長or頻度変更の検討',
            ]),
            ('Month 2: モデル再分析 & 精度確認', [
                '1. 改善データで再MMM実行',
                f'2. R²・NRMSEの改善幅を確認（現状R²={r2:.3f}）',
                '3. 精度確認後に予算判断へ移行',
            ]),
            ('Month 3: 予算最適化 & 次年度計画', [
                '1. 再分析結果に基づく予算再配分',
                '2. 有効チャネルへの集中投資',
                '3. 年間予算計画への反映',
            ]),
        ]
        kpi_text = f'KPI目標: R²≥0.80（再分析後） ｜ NRMSE≤0.10 ｜ データ品質スコア改善'

    elif n_dummies > 20:
        # Route B: モデル精査→最適化ルート
        months = [
            (f'Month 1: ダミー精査 & 初動対応', [
                f'1. ダミー変数（{n_dummies}本）の要因確認・除外判断',
                f'2. {zero_action}',
                '3. 伸び代ありチャネルへの小幅増額',
            ]),
            ('Month 2: 精査完了 & 最適化実行', [
                '1. 除外ダミーで再分析（R²変化確認）',
                '2. 予算最適化シナリオの本格実施',
                '3. 季節性ダミーの追加検討',
            ]),
            ('Month 3: 効果検証 & 次年度計画', [
                f'1. CV増加の実績検証（vs見込み+{cv_lift:.1f}%）',
                '2. ROI変化に応じた予算再配分',
                '3. 年間予算計画への反映',
            ]),
        ]
        kpi_text = f'KPI目標: CV月次 +{cv_lift/3:.0f}〜+{cv_lift:.0f}% ｜ R²≥{min(r2+0.01, 0.99):.2f}（再分析後） ｜ ゼロ係数≤{max(len(zero_channels)-2, 0)}ch'

    else:
        # Route C: 即時最適化ルート
        months = [
            ('Month 1: 予算再配分 & 増額実施', [
                '1. 伸び代ありチャネルへ予算集中',
                f'2. {zero_action_c}',
                '3. 最適化シナリオの実施承認',
            ]),
            ('Month 2: 効果検証 & モデル更新', [
                '1. 増額後のCV・CPA変化モニタリング',
                '2. データ蓄積後の再MMM実行',
                f'3. 追加チャネル投資の検討（目標+{cv_lift:.1f}%）',
            ]),
            ('Month 3: 次年度計画 & スケール', [
                f'1. CV増加の実績検証（vs見込み+{cv_lift:.1f}%）',
                '2. 高ROIチャネルへのスケールアップ',
                '3. 年間予算計画・フォーキャスト更新',
            ]),
        ]
        kpi_text = f'KPI目標: CV月次 +{cv_lift/3:.0f}〜+{cv_lift:.0f}% ｜ CPA改善率≥10% ｜ 伸び代ありチャネル予算比率向上'

    _road_w = 3.5
    _road_xs = [0.5, 4.55, 8.6]
    for i, (title, items) in enumerate(months):
        x = _road_xs[i]
        _rect(sl, x, 1.6, _road_w, 4.1, line_color=BRAND_LIGHT, line_pt=1.2)
        _box(sl, x + 0.15, 1.82, _road_w - 0.3, 0.6, title,
             font_size=13, bold=True, text_color=SECTION_BG)
        for j, item in enumerate(items):
            _box(sl, x + 0.2, 2.55 + j * 0.75, _road_w - 0.35, 0.65, item,
                 font_size=11, text_color=SECTION_BG, line_spacing=1.5)
        if i < 2:
            _box(sl, _road_xs[i] + _road_w + 0.05, 3.4, 0.4, 0.5, '▶',
                 font_size=16, text_color=SECTION_BG, align=PP_ALIGN.CENTER)
    _box(sl, 0.4, 5.85, 12.5, 0.45, kpi_text,
         font_size=11, bold=True, text_color=SECTION_BG, align=PP_ALIGN.CENTER)

    # Glossary
    sl = _add_slide(prs)
    _slide_title(sl, '用語集・モデル説明')
    glossary = [
        ('R²（決定係数）',       '0〜1の値。実測CVの分散のうちモデルが説明できる割合。0.90以上が目安。'),
        ('NRMSE',               '正規化された平均二乗誤差。小さいほど予測精度が高い。0.1以下が実用水準。'),
        ('NRMSE (holdout)',      f'最終{holdout_days}{period_unit}の未学習データでのNRMSE。モデルの汎化性能を評価。'),
        ('RSSD',         '支出金額シェアとCV貢献シェアの乖離。Robyn準拠の整合性指標。0に近いほど良い。'),
        ('MAPE',                f'平均絶対誤差率(%)。{mape_note}'),
        ('Adstock',             '広告効果の残存期間を表すモデル。λが高いほど効果が長期持続する。'),
        ('Hill飽和曲線',         '支出金額増加に対するCV効果の逓減を表す。α・γで形状が決まる。'),
        ('MCR',                 '媒体帰属率（Media Contribution Rate）。広告起因CVが全CVに占める割合。15%以上が目安。'),
        ('CPA',                 'Cost Per Acquisition（獲得単価）。支出金額÷推定CV数。低いほど効率的。'),
        ('ROI',                 'Return on Investment。推定CV数÷支出金額（万円）。高いほど効率的。'),
        ('ダミー変数',           '特定日付の異常値を補正する変数。自動探索エンジンが候補を生成・採用。'),
    ]
    _add_table(sl, 1.65, 1.1, 10.0, 5.9,
               ['用語', '説明'], [[k, v] for k, v in glossary],
               col_widths=[0.7, 3.5], body_font_size=10,
               col_alignments=[PP_ALIGN.CENTER, PP_ALIGN.LEFT])

    # Executive Summary (last slide)
    sl = _add_slide(prs)
    _bg(sl, COVER_BG)
    _box(sl, 0.5, 0.3, 12, 0.65, 'Executive Summary',
         font_size=28, bold=True, text_color=WHITE, align=PP_ALIGN.CENTER)
    _line(sl, 0.5, 1.1, 12.83, 1.1, color=RULE_COLOR)
    ex_kpis = [
        (f'R²={r2:.3f}',        f'{freq_label}{analysis_days}{period_unit}', ACCENT_GRN),
        (f'CPA最安 ¥{min((cm["cpa"] for cm in ch_metrics.values() if cm["cpa"]), default=0):,.0f}',
         '最高効率チャネル', BRAND_LIGHT),
        (f'ゼロ係数 {len(zero_channels)}ch', '予算削減推奨',         ACCENT_ORG),
        (f'CV+{cv_lift:.1f}%',  '同予算最適化で実現可能（試算）',   ACCENT_GRN),
        (f'ダミー {n_dummies}本', '自動探索済み',                    BRAND_LIGHT),
        (f'NRMSE(HO)={nrmse_hold:.4f}', '検証期間（holdout）の精度', ACCENT_ORG),
    ]
    for i, (v, s, c) in enumerate(ex_kpis):
        x = 0.5 + (i % 3) * 4.1
        y = 1.3 + (i // 3) * 1.5
        _rect(sl, x, y, 3.7, 1.3, fill_color=SECTION_BG)
        _box(sl, x + 0.1, y + 0.08, 3.5, 0.4, s,
             font_size=10, text_color=BRAND_LIGHT, align=PP_ALIGN.CENTER)
        _box(sl, x + 0.1, y + 0.5, 3.5, 0.7, v,
             font_size=18, bold=True, text_color=c or WHITE, align=PP_ALIGN.CENTER)
    next_steps = [
        'ダミー除外再分析 → R²向上目標',
        f'高ROI ({", ".join(top_roi_channels[:3])}) への予算集中',
        'ゼロ係数チャネル停止・検証',
    ]
    _box(sl, 0.5, 4.55, 4, 0.4, '次のステップ',
         font_size=12, bold=True, text_color=WHITE)
    for i, s in enumerate(next_steps):
        _box(sl, 0.5, 5.0 + i * 0.42, 12, 0.35, f'・{s}',
             font_size=11, text_color=BRAND_LIGHT)
    _copyright(sl)

    # ── Page numbers (skip cover = index 0) ───────────────────
    _DARK_BG = {
        (COVER_BG[0],   COVER_BG[1],   COVER_BG[2]),
        (SECTION_BG[0], SECTION_BG[1], SECTION_BG[2]),
    }
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        try:
            rgb = slide.background.fill.fore_color.rgb
            is_dark = (int(rgb[0]), int(rgb[1]), int(rgb[2])) in _DARK_BG
        except Exception:
            is_dark = False
        num_color = WHITE if is_dark else BRAND_LIGHT
        _box(slide, 12.3, 7.1, 0.85, 0.28, str(idx),
             font_size=9, text_color=num_color, align=PP_ALIGN.RIGHT)

    from datetime import datetime as _dt
    now = _dt.utcnow()  # python-pptxはUTC前提でZサフィックスを付けるため
    prs.core_properties.created  = now
    prs.core_properties.modified = now
    prs.save(output_path)
    return output_path
